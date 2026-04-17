"""Rich PPTX renderer with template bookend system.

When a template has ≥2 slides, the first and last are **preserved in-place**:
- **Cover** (first template slide): kept intact, title/subtitle filled into
  existing placeholders.  All original shapes/graphics preserved.
- **Closing** (last template slide): kept 100% untouched — all baked-in text,
  graphics, and design preserved exactly as the template provides.

Only middle example slides are deleted.  Content slides are inserted between
the kept bookends.  The closing slide's layout index is tracked via
``excluded_idx`` to prevent content slides from accidentally using it.
"""
from __future__ import annotations
import math
import re
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.oxml.ns import qn

from .schemas import (
    PresentationSpec, SlideSpec, SlideElement,
    TextContent, BulletContent, ChartContent, TableContent,
    ShapeContent, InfographicContent, SlideMasterInfo, ThemeColors,
)
from .slide_master import get_layout_for_slide_type, read_slide_master
from . import config
from .drawingml_effects import (
    add_shadow, add_gradient, add_theme_gradient,
    remove_outline, set_corner_radius, style_card,
    style_accent_bar, style_numbered_circle,
)
from .components import render_chart_container, render_accent_divider
from .color_utils import (
    FALLBACK_ACCENT_HEX, FLOW_ACCENT_HEX, CMP_ACCENT_HEX, KPI_ACCENT_HEX,
    pick_text_color, darken_hex, pick_text_color_for_brightness,
    effective_hex_after_brightness, abbreviate_number,
)

import logging
_log = logging.getLogger(__name__)

# Module-level slide dimensions — set at the start of render_presentation()
_sw: int = config.SLIDE_WIDTH
_sh: int = config.SLIDE_HEIGHT

# Placeholder types that must NEVER be removed (footers, slide numbers, dates)
_PROTECTED_PH_TYPES = frozenset({
    12,  # SLIDE_NUMBER
    11,  # DATE_TIME
    10,  # FOOTER
})


# ── Chart type mapping ──
CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

# ── Shape type mapping ──
SHAPE_MAP = {
    "ROUNDED_RECTANGLE": MSO_SHAPE.ROUNDED_RECTANGLE,
    "RECTANGLE": MSO_SHAPE.RECTANGLE,
    "CHEVRON": MSO_SHAPE.CHEVRON,
    "RIGHT_ARROW": MSO_SHAPE.RIGHT_ARROW,
    "OVAL": MSO_SHAPE.OVAL,
    "PENTAGON": MSO_SHAPE.PENTAGON,
    "HEXAGON": MSO_SHAPE.HEXAGON,
    "DIAMOND": MSO_SHAPE.DIAMOND,
    "FLOWCHART_PROCESS": MSO_SHAPE.FLOWCHART_PROCESS,
    "FLOWCHART_DECISION": MSO_SHAPE.FLOWCHART_DECISION,
    "FLOWCHART_TERMINATOR": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "ROUND_1_RECTANGLE": MSO_SHAPE.ROUND_1_RECTANGLE,
    "SNIP_1_RECTANGLE": MSO_SHAPE.SNIP_1_RECTANGLE,
}


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex string like '1F77B4' to RGBColor."""
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ── Theme color indices for shape fills (auto-inherit Slide Master palette) ──
_ACCENT_THEME_COLORS = [
    MSO_THEME_COLOR.ACCENT_1, MSO_THEME_COLOR.ACCENT_2,
    MSO_THEME_COLOR.ACCENT_3, MSO_THEME_COLOR.ACCENT_4,
    MSO_THEME_COLOR.ACCENT_5, MSO_THEME_COLOR.ACCENT_6,
]

# Fallback hex palette — WCAG-compliant (≥4.5:1 vs white text)
_FALLBACK_ACCENT_HEX = FALLBACK_ACCENT_HEX


def _apply_accent_fill(shape, index: int, has_template: bool) -> None:
    """Apply an accent color fill to a shape, using theme colors when a template is loaded."""
    shape.fill.solid()
    if has_template:
        shape.fill.fore_color.theme_color = _ACCENT_THEME_COLORS[index % len(_ACCENT_THEME_COLORS)]
    else:
        shape.fill.fore_color.rgb = _hex_to_rgb(_FALLBACK_ACCENT_HEX[index % len(_FALLBACK_ACCENT_HEX)])


def _get_accent_hex(master_info: SlideMasterInfo | None, index: int) -> str:
    """Get accent hex color from theme or fallback."""
    if master_info:
        accents = master_info.theme_colors.accents()
        return accents[index % len(accents)]
    return _FALLBACK_ACCENT_HEX[index % len(_FALLBACK_ACCENT_HEX)]


def _remove_unused_placeholders(slide) -> None:
    """Remove placeholder shapes that have no user-supplied text.

    Preserves footer, slide-number, and date-time placeholders so that
    the template's built-in furniture is not destroyed.
    """
    for shape in list(slide.placeholders):
        ph_idx = shape.placeholder_format.idx
        ph_type = shape.placeholder_format.type  # int enum
        # Never remove protected placeholders
        if int(ph_type) in _PROTECTED_PH_TYPES or ph_idx in (10, 11, 12):
            continue
        # Also skip any placeholder in the bottom 15% (likely footer zone)
        if hasattr(shape, 'top') and shape.top > _sh * 0.85:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip() if shape.text_frame.text else ""
            if not text:
                try:
                    sp = shape._element
                    sp.getparent().remove(sp)
                except Exception:
                    pass


def _set_autofit(text_frame, *, shrink_ok: bool = False) -> None:
    """Set auto-size mode for a text frame.

    *shrink_ok=True* enables TEXT_TO_FIT_SHAPE for bounded shapes (cards,
    badges) so that long text shrinks instead of overflowing.
    """
    text_frame.auto_size = (
        MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if shrink_ok else MSO_AUTO_SIZE.NONE
    )
    text_frame.word_wrap = True


def _set_text_frame_text(text_frame, text: str, font_size=None, bold: bool | None = None,
                         alignment=None, color_rgb: RGBColor | None = None,
                         theme_color: MSO_THEME_COLOR | None = None,
                         font_name: str | None = None) -> None:
    """Replace a text frame with a single formatted paragraph."""
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = config.TF_MARGIN_LEFT
    text_frame.margin_right = config.TF_MARGIN_RIGHT
    text_frame.margin_top = config.TF_MARGIN_TOP
    text_frame.margin_bottom = config.TF_MARGIN_BOTTOM
    p = text_frame.paragraphs[0]
    p.text = text
    if font_size is not None:
        p.font.size = font_size
    if bold is not None:
        p.font.bold = bold
    if alignment is not None:
        p.alignment = alignment
    if theme_color is not None:
        p.font.color.theme_color = theme_color
    elif color_rgb is not None:
        p.font.color.rgb = color_rgb
    if font_name:
        p.font.name = font_name
    _set_autofit(text_frame)


def _populate_text_list(text_frame, items: list[str], font_size, prefix: str = "",
                        font_name: str | None = None) -> None:
    """Populate a text frame with a concise multi-paragraph list."""
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = config.TF_MARGIN_LEFT
    text_frame.margin_right = config.TF_MARGIN_RIGHT
    text_frame.margin_top = config.TF_MARGIN_TOP
    text_frame.margin_bottom = config.TF_MARGIN_BOTTOM
    for idx, item in enumerate(items):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = f"{prefix}{item}" if prefix else item
        p.font.size = font_size
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = Pt(int(font_size.pt * 1.4)) if hasattr(font_size, 'pt') else None
        if idx > 0:
            p.space_before = config.BULLET_SPACE_BEFORE
        if font_name:
            p.font.name = font_name
    _set_autofit(text_frame, shrink_ok=True)


def _apply_light_surface_fill(shape, has_tpl: bool, brightness: float = 0.96,
                              fallback_hex: str = "F7F9FC") -> None:
    """Apply a subtle, low-noise surface fill for professional content containers."""
    shape.fill.solid()
    if has_tpl:
        shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        shape.fill.fore_color.brightness = brightness
    else:
        shape.fill.fore_color.rgb = _hex_to_rgb(fallback_hex)
    shape.line.fill.background()


def _remove_text_artifacts(slide) -> None:
    """Remove non-placeholder text shapes inherited from a layout to avoid ghost copy."""
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False):
            continue
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip() if shape.text_frame.text else ""
            if text:
                try:
                    sp = shape._element
                    sp.getparent().remove(sp)
                except Exception:
                    pass


def _render_slide_title(slide, title: str, subtitle: str | None = None, has_tpl: bool = False) -> None:
    """Render title using the template title placeholder when available."""
    title_shape = slide.shapes.title if has_tpl else None
    if title_shape is not None and title_shape.has_text_frame:
        _set_text_frame_text(
            title_shape.text_frame,
            title,
            font_size=config.FONT_TITLE,
            bold=True,
            alignment=PP_ALIGN.LEFT,
            font_name=config.FONT_NAME_PRIMARY,
        )
        return
    _add_title_bar(slide, title, subtitle, has_tpl)


def _apply_brand_card_fill(shape, index: int, has_tpl: bool,
                           master_info: SlideMasterInfo | None = None) -> str:
    """Apply brand-aligned card background: ACCENT_1 with varying brightness.

    Returns the *effective* hex color after brightness — callers should use
    this to pick contrast-correct text colors instead of hardcoding white.
    """
    brightness = config.CARD_BRIGHTNESS_LEVELS[index % len(config.CARD_BRIGHTNESS_LEVELS)]
    shape.fill.solid()
    if has_tpl:
        shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        shape.fill.fore_color.brightness = brightness
        base_hex = (
            master_info.theme_colors.accents()[0]
            if master_info and master_info.theme_colors.accents()
            else FALLBACK_ACCENT_HEX[0]
        )
    else:
        base_hex = "EDF2F9"
        shape.fill.fore_color.rgb = _hex_to_rgb(base_hex)
    return effective_hex_after_brightness(base_hex, brightness)


def _add_speaker_notes(slide, text: str) -> None:
    """Add speaker notes to a slide."""
    try:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text[:2000]
    except Exception:
        pass  # notes slide creation can fail on some templates


def render_presentation(spec: PresentationSpec, output_path: str | Path) -> Path:
    """Render a full presentation from PresentationSpec to .pptx file."""
    global _sw, _sh
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    template_path = spec.template_path
    last_slide_layout_idx = None

    if template_path and Path(template_path).exists():
        # Load template to get theme colors/layouts
        prs = Presentation(str(template_path))
        n_tpl_slides = len(prs.slides)
        
        # Enable bookend system if template has ≥2 slides (preserve first and last)
        if n_tpl_slides >= 2:
            template_has_bookends = True
            # Capture layout index of last slide to prevent content slides from using it
            # Find the layout index manually since SlideLayout objects don't have an .index attribute
            closing_layout = prs.slides[n_tpl_slides - 1].slide_layout
            last_slide_layout_idx = -1
            for idx, layout in enumerate(prs.slide_layouts):
                if layout == closing_layout:
                    last_slide_layout_idx = idx
                    break
            
            # Delete only middle slides (indices 1 to n-2), preserve first and last
            for i in range(n_tpl_slides - 2, 0, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])
            _log.info(f"Template loaded with {n_tpl_slides} slides - preserved first and last as bookends")
        else:
            # Template has <2 slides, delete all and don't use bookend system
            template_has_bookends = False
            last_slide_layout_idx = None
            for i in range(n_tpl_slides - 1, -1, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])
            _log.info(f"Template loaded with {n_tpl_slides} slides - deleted all (insufficient for bookends)")
    else:
        prs = Presentation()
        prs.slide_width = config.SLIDE_WIDTH
        prs.slide_height = config.SLIDE_HEIGHT
        template_has_bookends = False
        last_slide_layout_idx = None

    # Capture actual slide dimensions for the rest of the render chain
    _sw = int(prs.slide_width)
    _sh = int(prs.slide_height)
    _log.info(f"Slide dimensions: {_sw}x{_sh} EMU ({_sw/914400:.2f}x{_sh/914400:.2f} in)")

    master_info = read_slide_master(template_path) if template_path and Path(template_path).exists() else None

    if template_has_bookends:
        # Cover slide (first template slide) — fill title/subtitle into placeholders
        cover_slide = prs.slides[0]
        phs = sorted(cover_slide.placeholders, key=lambda p: p.placeholder_format.idx)
        if len(phs) >= 1 and spec.title:
            phs[0].text = spec.title
            _set_autofit(phs[0].text_frame)
        if len(phs) >= 2 and spec.subtitle:
            phs[1].text = spec.subtitle
            _set_autofit(phs[1].text_frame)

        # Presenter Name + Date metadata (mimics reference samples)
        _add_cover_metadata(
            cover_slide,
            presenter=getattr(spec, "presenter", "") or "",
            date_str=getattr(spec, "date_str", "") or "",
            master_info=master_info,
            has_tpl=True,
        )

        # Closing slide (last template slide) — 100% untouched.
        # It currently sits at index 1 (after middle deletions).
        #
        # IMPORTANT: we must NOT detach its sldId from sldIdLst before adding
        # new content slides.  python-pptx's CT_SlideIdList._next_id computes
        # the next slide id as ``max(existing_ids) + 1`` from the *current*
        # sldIdLst.  If closing (id=280 for our templates) is detached, its id
        # becomes invisible, causing new slides to be assigned the same id and
        # producing a duplicate ``p:sldId/@id`` when closing is re-inserted.
        # PowerPoint rejects such files as corrupt ("can't read this").
        #
        # Instead, we keep closing in place while rendering so its id is
        # visible to _next_id (new slides start at closing_id + 1), then move
        # it to the end.  lxml's ``append`` on an element that is already in
        # the tree *moves* it atomically (remove-then-append inside the parent)
        # so no duplicate is ever created.
        _closing_sldId = prs.slides._sldIdLst[1]

        # Render content slides (skip cover and thank_you — template provides them).
        # New sldId elements are appended AFTER the closing entry.
        presenter = getattr(spec, "presenter", "") or ""
        date_str = getattr(spec, "date_str", "") or ""
        for slide_spec in spec.slides:
            if slide_spec.slide_type in ("cover", "thank_you"):
                continue
            _render_slide(prs, slide_spec, master_info, deck_title=spec.title,
                          excluded_layout_idx=last_slide_layout_idx,
                          presenter=presenter, date_str=date_str)

        # Move closing slide from index 1 to the end.  append() on an element
        # already parented will detach and re-insert it — no duplicate.
        prs.slides._sldIdLst.append(_closing_sldId)

        # Renumber slide part names to avoid duplicate zip entry warnings
        rIds = [sldId.rId for sldId in prs.slides._sldIdLst]
        prs.part.rename_slide_parts(rIds)
    else:
        presenter = getattr(spec, "presenter", "") or ""
        date_str = getattr(spec, "date_str", "") or ""
        for slide_spec in spec.slides:
            _render_slide(prs, slide_spec, master_info, deck_title=spec.title,
                          presenter=presenter, date_str=date_str)

    prs.save(str(output_path))
    return output_path




def _get_slide_layout(prs: Presentation, slide_spec: SlideSpec, master_info: SlideMasterInfo | None,
                      excluded_layout_idx: int | None = None):
    """Get the appropriate slide layout from the presentation."""
    if master_info:
        layout_info = get_layout_for_slide_type(master_info, slide_spec.slide_type,
                                                excluded_idx=excluded_layout_idx)
        # Find the actual layout object by index
        master = prs.slide_masters[0]
        if layout_info.index < len(master.slide_layouts):
            return master.slide_layouts[layout_info.index]

    # Fallback: use layout by name matching or index
    for layout in prs.slide_layouts:
        if slide_spec.layout_name.lower() in layout.name.lower():
            return layout

    # Last fallback: blank-like layout (usually index 6 for blank in default)
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout

    return prs.slide_layouts[0]


def _render_slide(prs: Presentation, spec: SlideSpec, master_info: SlideMasterInfo | None,
                  deck_title: str = "", excluded_layout_idx: int | None = None,
                  presenter: str = "", date_str: str = ""):
    """Render a single slide."""
    layout = _get_slide_layout(prs, spec, master_info, excluded_layout_idx=excluded_layout_idx)
    slide = prs.slides.add_slide(layout)
    has_tpl = master_info is not None

    # ── Cover slide ──
    if spec.slide_type == "cover":
        _render_cover(slide, spec, master_info, has_tpl)
        _add_cover_metadata(slide, presenter, date_str, master_info, has_tpl)
        _add_speaker_notes(slide, f"Cover: {spec.title}")
        return

    # ── Thank you slide ──
    if spec.slide_type == "thank_you":
        _render_thank_you(slide, spec, master_info, has_tpl)
        _add_speaker_notes(slide, "Thank you slide — open for questions and discussion.")
        return

    # ── Section divider ──
    if spec.slide_type == "section_divider":
        _render_divider(slide, spec, has_tpl)
        _remove_unused_placeholders(slide)
        _add_speaker_notes(slide, f"Section: {spec.title}" + (f" — {spec.subtitle}" if spec.subtitle else ""))
        return

    # ── All other slides: populate title first, then add content, then clean placeholders ──
    if spec.title:
        _render_slide_title(slide, spec.title, spec.subtitle, has_tpl)

    for element in spec.elements:
        try:
            _render_element(slide, element, master_info, has_tpl, spec.slide_type)
        except Exception as e:
            _log.warning(f"Skipped element {element.element_type} on slide {spec.slide_number}: {e}")

    _remove_unused_placeholders(slide)

    # ── Slide furniture (footer, accent stripe) ──
    _add_slide_furniture(slide, spec, has_tpl, deck_title, master_info)

    # ── Speaker notes ──
    notes_text = spec.title or ""
    if spec.subtitle:
        notes_text += f"\n{spec.subtitle}"
    _add_speaker_notes(slide, notes_text)


# ── Slide type renderers ────────────────────────────────────────────

def _shape_intersects_metadata_zone(shape, zone_top: int, zone_bottom: int,
                                    zone_left: int, zone_right: int) -> bool:
    """Check if a shape's bounding box overlaps the metadata insertion zone."""
    try:
        s_left = int(shape.left or 0)
        s_top = int(shape.top or 0)
        s_right = s_left + int(shape.width or 0)
        s_bottom = s_top + int(shape.height or 0)
    except Exception:
        return False
    # Axis-aligned bounding-box intersection test
    if s_right < zone_left or s_left > zone_right:
        return False
    if s_bottom < zone_top or s_top > zone_bottom:
        return False
    return True


def _cover_text_color_rgb(slide) -> RGBColor:
    """Return the best text colour for cover-slide metadata (presenter/date).

    Scans the cover slide + its layout for a *concrete* dark signal:
    - A large solid-fill rectangle/freeform with low luminance, OR
    - A fullscreen picture *combined with* a dark overlay shape.

    Pictures alone are ambiguous (the Accenture cover uses a pink picture
    that needs dark text, the UAE cover uses a green-overlay photo that
    needs white) so we require an explicit dark-coloured overlay before
    flipping to white. Defaults to the near-black body colour otherwise.
    """
    from .color_utils import relative_luminance_rgb
    try:
        slide_area = float(_sw) * float(_sh)
        dark_fill_area = 0.0
        # Inspect both the slide's own shapes and its layout's shapes — most
        # cover artwork lives on the layout, not the slide itself.
        shape_pools = [slide.shapes]
        try:
            shape_pools.append(slide.slide_layout.shapes)
        except Exception:
            pass
        for pool in shape_pools:
            for sh in pool:
                if not (getattr(sh, "width", None) and getattr(sh, "height", None)):
                    continue
                area = float(sh.width) * float(sh.height)
                if area < slide_area * 0.25:
                    continue
                # Only solid-fill shapes count as a dark-overlay signal.
                try:
                    fill = sh.fill
                    if fill.type == 1:  # SOLID
                        rgb = fill.fore_color.rgb
                        if rgb is not None and relative_luminance_rgb(rgb) < 0.35:
                            dark_fill_area += area
                except Exception:
                    continue
        if dark_fill_area >= slide_area * 0.5:
            return RGBColor(0xFF, 0xFF, 0xFF)
    except Exception:
        pass
    return RGBColor(0x22, 0x22, 0x22)


def _add_cover_metadata(slide, presenter: str, date_str: str,
                        master_info: SlideMasterInfo | None = None,
                        has_tpl: bool = False) -> None:
    """Add ``Presenter Name: ...`` and ``Date: ...`` textboxes to the cover slide.

    Positions match the reference Caspr + Ghost Research samples:
    ``x ≈ 0.375 in`` (~342900 EMU), bottom-left. Skips if either (a) both
    values are empty, or (b) existing cover artwork already covers that zone.
    """
    presenter = (presenter or "").strip()
    date_str = (date_str or "").strip()
    if not presenter and not date_str:
        return

    # Resolve a contrast-appropriate text colour for the cover backdrop.
    text_color_rgb = _cover_text_color_rgb(slide)

    # Insertion zone: ~0.375in from left, spanning from 6.45in to 7.10in vertical
    zone_left = Emu(342900)
    zone_right = Emu(4_400_000)   # ~4.8 in wide
    zone_top = Emu(5_800_000)
    zone_bottom = Emu(6_500_000)

    # Collision check against existing shapes
    for sh in list(slide.shapes):
        try:
            if _shape_intersects_metadata_zone(sh, int(zone_top), int(zone_bottom),
                                                int(zone_left), int(zone_right)):
                # Existing artwork or placeholder occupies the zone — don't double up
                _log.debug("Skipping cover metadata: zone already occupied by %r", sh.name)
                return
        except Exception:
            continue

    # Text style: bold, small, theme-dark text so it reads as an editorial tag
    box_w = Emu(4_000_000)
    box_h = Emu(280_000)

    if presenter:
        presenter_box = slide.shapes.add_textbox(
            zone_left, Emu(5_905_418), box_w, box_h,
        )
        tf = presenter_box.text_frame
        tf.word_wrap = False
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.text = f"Presenter Name: {presenter}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = config.FONT_NAME_PRIMARY
        p.font.color.rgb = text_color_rgb
        p.alignment = PP_ALIGN.LEFT

    if date_str:
        date_box = slide.shapes.add_textbox(
            zone_left, Emu(6_244_959), box_w, box_h,
        )
        tf = date_box.text_frame
        tf.word_wrap = False
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.text = f"Date: {date_str}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = config.FONT_NAME_PRIMARY
        p.font.color.rgb = text_color_rgb
        p.alignment = PP_ALIGN.LEFT


def _render_cover(slide, spec: SlideSpec, master_info: SlideMasterInfo | None = None,
                  has_tpl: bool = False):
    """Render cover slide using placeholders or fallback to shapes."""
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    if phs:
        ph_list = sorted(phs.values(), key=lambda p: p.placeholder_format.idx)
        if len(ph_list) >= 1:
            ph_list[0].text = spec.title
            for para in ph_list[0].text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = config.FONT_TITLE
                    run.font.name = config.FONT_NAME_PRIMARY
        if len(ph_list) >= 2 and spec.subtitle:
            ph_list[1].text = spec.subtitle
            for para in ph_list[1].text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = config.FONT_SUBTITLE
                    run.font.name = config.FONT_NAME_PRIMARY
        # Remove remaining unused placeholders to prevent ghost text
        for ph in ph_list[2:]:
            if not ph.text.strip():
                try:
                    sp = ph._element
                    sp.getparent().remove(sp)
                except Exception:
                    ph.text = ""  # fallback: blank it out
        # Auto-fit title/subtitle text
        if len(ph_list) >= 1:
            _set_autofit(ph_list[0].text_frame)
        if len(ph_list) >= 2:
            _set_autofit(ph_list[1].text_frame)
    else:
        cw = _sw - int(config.MARGIN_LEFT) - int(config.MARGIN_RIGHT)
        _add_textbox(slide, spec.title, config.MARGIN_LEFT, Emu(2400000),
                     cw, Emu(900000),
                     font_size=config.FONT_TITLE, bold=True, alignment="center",
                     font_name=config.FONT_NAME_PRIMARY)
        if spec.subtitle:
            _add_textbox(slide, spec.subtitle, config.MARGIN_LEFT, Emu(3500000),
                         cw, Emu(600000),
                         font_size=config.FONT_SUBTITLE, alignment="center",
                         font_name=config.FONT_NAME_PRIMARY)

    # Bottom accent bar on cover with gradient
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, _sh - Emu(100000),
        _sw, Emu(100000)
    )
    if has_tpl:
        style_accent_bar(accent, theme_color=MSO_THEME_COLOR.ACCENT_1, angle=0.0)
    else:
        add_gradient(accent, [(0.0, _hex_to_rgb("4472C4")), (1.0, _hex_to_rgb("2B579A"))], angle=0.0)
        remove_outline(accent)

    for element in spec.elements:
        _render_element(slide, element, master_info, has_tpl)


def _render_divider(slide, spec: SlideSpec, has_tpl: bool = False):
    """Render a section divider slide with title and optional subtitle."""
    cw = _sw - int(config.MARGIN_LEFT) - int(config.MARGIN_RIGHT)
    subtitle_handled = False
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if phs:
        ph_list = sorted(phs.values(), key=lambda p: p.placeholder_format.idx)
        if ph_list:
            ph_list[0].text = spec.title
            _set_autofit(ph_list[0].text_frame)
        # Fill subtitle into second placeholder or remove it
        if len(ph_list) >= 2:
            if spec.subtitle:
                ph_list[1].text = spec.subtitle
                _set_autofit(ph_list[1].text_frame)
                subtitle_handled = True
            else:
                try:
                    sp = ph_list[1]._element
                    sp.getparent().remove(sp)
                except Exception:
                    ph_list[1].text = ""
        # Remove any remaining unused placeholders (except protected)
        for ph in ph_list[2:]:
            ph_type = ph.placeholder_format.type
            if int(ph_type) in _PROTECTED_PH_TYPES or ph.placeholder_format.idx in (10, 11, 12):
                continue
            try:
                sp = ph._element
                sp.getparent().remove(sp)
            except Exception:
                ph.text = ""
    else:
        _add_textbox(slide, spec.title, config.MARGIN_LEFT, Emu(2700000),
                     cw, Emu(900000),
                     font_size=config.FONT_TITLE, bold=True, alignment="center",
                     font_name=config.FONT_NAME_PRIMARY)

    # Add subtitle as manual textbox when placeholders didn't handle it
    if spec.subtitle and not subtitle_handled:
        _add_textbox(slide, spec.subtitle, config.MARGIN_LEFT, Emu(3700000),
                     cw, Emu(500000),
                     font_size=config.FONT_SUBTITLE, alignment="center", color="666666",
                     font_name=config.FONT_NAME_PRIMARY)

    # Accent bar under title with gradient — center it horizontally
    bar_y = Emu(4300000) if spec.subtitle else Emu(3800000)
    bar_w = min(Emu(5000000), cw)
    bar_x = int(config.MARGIN_LEFT) + (cw - bar_w) // 2
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        bar_x, bar_y,
        bar_w, Emu(36000)
    )
    if has_tpl:
        style_accent_bar(bar, theme_color=MSO_THEME_COLOR.ACCENT_1, angle=0.0)
    else:
        add_gradient(bar, [(0.0, _hex_to_rgb("4472C4")), (1.0, _hex_to_rgb("2B579A"))], angle=0.0)
        remove_outline(bar)


def _render_thank_you(slide, spec: SlideSpec, master_info: SlideMasterInfo | None = None,
                      has_tpl: bool = False):
    """Render thank you slide — use EITHER placeholders OR textboxes, never both.

    When a template is loaded the layout already contains styled "Thank You"
    text, so we leave the slide completely untouched to preserve the design.
    """
    if has_tpl:
        # Template thank-you layouts usually have the text baked into layout
        # shapes. However, some layouts have 0 placeholders and text may not
        # be visible. Add a subtle insurance textbox if the slide is empty.
        if len(list(slide.shapes)) == 0:
            cw = _sw - int(config.MARGIN_LEFT) - int(config.MARGIN_RIGHT)
            inset = min(Emu(600000), cw // 10)
            _add_textbox(slide, spec.title or "Thank You",
                         int(config.MARGIN_LEFT) + inset, Emu(2800000),
                         cw - 2 * inset, Emu(800000),
                         font_size=Pt(32), bold=True, alignment="center",
                         font_name=config.FONT_NAME_PRIMARY)
        return

    # No template — render manually
    title_text = spec.title or "Thank You"
    subtitle_text = spec.subtitle or "Questions & Discussion"

    _remove_text_artifacts(slide)
    cw = _sw - int(config.MARGIN_LEFT) - int(config.MARGIN_RIGHT)
    inset = min(Emu(600000), cw // 10)
    _add_textbox(slide, title_text, int(config.MARGIN_LEFT) + inset, Emu(2300000),
                 cw - 2 * inset, Emu(900000),
                 font_size=Pt(36), bold=True, alignment="center",
                 font_name=config.FONT_NAME_PRIMARY)
    _add_textbox(slide, subtitle_text, int(config.MARGIN_LEFT) + inset, Emu(3300000),
                 cw - 2 * inset, Emu(400000),
                 font_size=config.FONT_SUBTITLE, alignment="center",
                 color="666666", font_name=config.FONT_NAME_PRIMARY)


# ── Title bar ───────────────────────────────────────────────────────

def _add_title_bar(slide, title: str, subtitle: str | None = None, has_tpl: bool = False):
    """Add a title bar at the top of a content slide with accent underline."""
    cw = _sw - int(config.MARGIN_LEFT) - int(config.MARGIN_RIGHT)
    # Title
    txBox = slide.shapes.add_textbox(
        config.MARGIN_LEFT, config.MARGIN_TOP,
        cw, Emu(530000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = config.FONT_TITLE
    p.font.bold = True
    p.font.name = config.FONT_NAME_PRIMARY
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            config.MARGIN_LEFT, Emu(config.MARGIN_TOP + 580000),
            cw, Emu(350000)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = Emu(0)
        tf2.margin_top = Emu(0)
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = config.FONT_SUBTITLE
        p2.font.name = config.FONT_NAME_PRIMARY
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p2.alignment = PP_ALIGN.LEFT


# ── Slide furniture (footer bar, accent stripe) ──────────────────────

def _pick_readable_accent_hex(master_info: SlideMasterInfo | None) -> str:
    """Return the first accent hex that has enough contrast on white.

    Some templates (e.g. UAE Solar) place a very light cream as ACCENT_1
    which renders the footer page number nearly invisible. This helper
    walks the available accents and picks the first that meets a ~3:1
    contrast ratio against a white slide background, falling back to a
    safe dark grey if none qualify.
    """
    from .color_utils import contrast_ratio, FALLBACK_ACCENT_HEX
    if not master_info or not master_info.theme_colors.accents():
        return "4472C4"
    for hex_val in master_info.theme_colors.accents():
        if not hex_val:
            continue
        if contrast_ratio(hex_val, "FFFFFF") >= 3.0:
            return hex_val.lstrip("#")
    return FALLBACK_ACCENT_HEX


def _add_slide_furniture(slide, spec: SlideSpec, has_tpl: bool, deck_title: str,
                         master_info: SlideMasterInfo | None = None):
    """Add the minimalist bottom rail + accent stripe + decorative chevrons.

    The 2024-style reference decks (Ghost Research / Caspr) use a *thin* bottom
    rail — a single hairline accent rule with a tiny dot and a right-aligned
    page number, not a heavy black footer bar. This function now emits that
    cleaner furniture so our slides read as editorial rather than corporate.

    *master_info* is used to pick an accent colour with enough contrast for
    the page number; templates whose ACCENT_1 is a very light tint would
    otherwise render the number invisible on a white slide.
    """
    accent_hex = _pick_readable_accent_hex(master_info)
    accent_rgb = _hex_to_rgb(accent_hex)

    # Thin strip at the bottom for the deck title + page number. No filled
    # bar — just a typographic rail so the footer feels light.
    footer_h = Emu(340000)
    footer_y = _sh - footer_h
    rail_y = footer_y + Emu(40000)

    # Hairline separator above the rail (0.75pt, readable accent)
    try:
        sep_line = slide.shapes.add_connector(
            1,
            config.MARGIN_LEFT, rail_y,
            Emu(_sw - int(config.MARGIN_RIGHT)), rail_y,
        )
        sep_line.line.color.rgb = accent_rgb
        sep_line.line.width = Pt(0.75)
    except Exception:
        pass

    # Small accent dot on the far-left of the rail (visual anchor)
    dot_d = Emu(110000)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        config.MARGIN_LEFT, rail_y - dot_d // 2,
        dot_d, dot_d,
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent_rgb
    remove_outline(dot)

    # Deck-title caption (muted grey, left-aligned, under the rail)
    if deck_title:
        ft = slide.shapes.add_textbox(
            config.MARGIN_LEFT + Emu(200000),
            rail_y + Emu(40000),
            Emu(min(8500000, _sw - 1800000)),
            Emu(230000),
        )
        tf = ft.text_frame
        tf.word_wrap = False
        tf.margin_left = Emu(0)
        tf.margin_top = Emu(0)
        p = tf.paragraphs[0]
        p.text = deck_title[:90]
        p.font.size = Pt(8)
        p.font.name = config.FONT_NAME_PRIMARY
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.alignment = PP_ALIGN.LEFT

    # Slide number (right-aligned, readable accent, bold)
    sn = slide.shapes.add_textbox(
        Emu(_sw - int(config.MARGIN_RIGHT) - 600000),
        rail_y + Emu(40000),
        Emu(600000), Emu(230000),
    )
    tf2 = sn.text_frame
    tf2.word_wrap = False
    tf2.margin_left = Emu(0)
    tf2.margin_right = Emu(0)
    tf2.margin_top = Emu(0)
    p2 = tf2.paragraphs[0]
    p2.text = str(spec.slide_number)
    p2.font.size = Pt(10)
    p2.font.name = config.FONT_NAME_PRIMARY
    p2.font.bold = True
    p2.font.color.rgb = accent_rgb
    p2.alignment = PP_ALIGN.RIGHT

    # ── 2. Left accent stripe (thin vertical bar) ──
    stripe_w = Emu(50000) if has_tpl else Emu(60000)
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, stripe_w, footer_y
    )
    if has_tpl:
        style_accent_bar(stripe, theme_color=MSO_THEME_COLOR.ACCENT_1, angle=270.0)
    else:
        add_gradient(stripe, [(0.0, _hex_to_rgb("4472C4")), (1.0, _hex_to_rgb("2B579A"))], angle=270.0)
        remove_outline(stripe)

    # ── 3. Title accent underline ──
    underline_y = Emu(1220000)
    underline_w = Emu(2200000)
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        config.MARGIN_LEFT, underline_y,
        underline_w, Emu(36000)
    )
    if has_tpl:
        style_accent_bar(underline, theme_color=MSO_THEME_COLOR.ACCENT_1, angle=0.0)
    else:
        add_gradient(underline, [(0.0, _hex_to_rgb("4472C4")), (1.0, _hex_to_rgb("2B579A"))], angle=0.0)
        remove_outline(underline)

    # ── 4. Top-right decorative chevron ──
    chev_size = Emu(180000)
    chev = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        _sw - Emu(600000), Emu(200000),
        chev_size, chev_size
    )
    if has_tpl:
        chev.fill.solid()
        chev.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        chev.fill.fore_color.brightness = 0.6
    else:
        chev.fill.solid()
        chev.fill.fore_color.rgb = _hex_to_rgb("D6E4F0")
    chev.line.fill.background()

    # Second chevron (slightly offset for layered effect)
    chev2 = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        _sw - Emu(420000), Emu(200000),
        chev_size, chev_size
    )
    if has_tpl:
        chev2.fill.solid()
        chev2.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        chev2.fill.fore_color.brightness = 0.8
    else:
        chev2.fill.solid()
        chev2.fill.fore_color.rgb = _hex_to_rgb("E8F0FA")
    chev2.line.fill.background()

    # NOTE: prior versions also drew a heavy 1.5pt rule + decorative circle
    # immediately above the footer. Both were duplicating the hairline rail
    # we now render at the top of this function and have been removed.


# ── Element renderer dispatch ───────────────────────────────────────

def _render_element(slide, element: SlideElement,
                    master_info: SlideMasterInfo | None = None,
                    has_tpl: bool = False,
                    slide_type: str = "content"):
    """Dispatch rendering based on element type."""
    pos = element.position
    content = element.content

    if isinstance(content, TextContent):
        _render_text(slide, pos, content)
    elif isinstance(content, BulletContent):
        _render_bullets(slide, pos, content, has_tpl, slide_type, master_info)
    elif isinstance(content, ChartContent):
        _render_chart(slide, pos, content, master_info, has_tpl)
    elif isinstance(content, TableContent):
        _render_table(slide, pos, content, master_info, has_tpl)
    elif isinstance(content, ShapeContent):
        _render_shape(slide, pos, content)
    elif isinstance(content, InfographicContent):
        _render_infographic(slide, pos, content, master_info, has_tpl)


# ── Text rendering ──────────────────────────────────────────────────

# Simple code detection for bullet items
def _looks_like_code(text: str) -> bool:
    """Quick heuristic to detect code-like content for monospace font."""
    if not text:
        return False
    # Patterns: backticks, YAML key:value, JSON brackets, code keywords
    patterns = [
        r'\`[^`]+\`',  # inline code
        r'^\s*[\w\-]+:\s*\S+',  # YAML-like
        r'[{\[\]}]',  # JSON brackets
        r'\b(def|class|function|const|let|var)\b',  # code keywords
    ]
    matches = sum(1 for p in patterns if re.search(p, text, re.MULTILINE))
    return matches >= 1 or text.count('`') >= 2


def _add_textbox(slide, text, left, top, width, height,
                 font_size=None, bold=False, italic=False,
                 color=None, alignment=None, autofit=True, font_name=None):
    """Helper to add a text box with formatting and optional auto-fit."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    p = tf.paragraphs[0]
    p.text = text
    if font_size:
        p.font.size = font_size
    if bold:
        p.font.bold = True
    if italic:
        p.font.italic = True
    if color:
        p.font.color.rgb = _hex_to_rgb(color) if isinstance(color, str) else color
    if alignment and alignment in ALIGN_MAP:
        p.alignment = ALIGN_MAP[alignment]
    if font_name:
        p.font.name = font_name
    if autofit:
        _set_autofit(tf)
    else:
        tf.auto_size = None
    return txBox


def _render_text(slide, pos, content: TextContent):
    """Render a text element."""
    # Determine font: use content's font_name or default to Inter
    font_name = content.font_name or config.FONT_NAME_PRIMARY
    _add_textbox(
        slide, content.text,
        pos.left, pos.top, pos.width, pos.height,
        font_size=Pt(content.font_size) if content.font_size else config.FONT_BODY,
        bold=content.bold, italic=content.italic,
        color=content.color, alignment=content.alignment,
        font_name=font_name,
    )


def _render_bullets(slide, pos, content: BulletContent, has_tpl: bool = False,
                    slide_type: str = "content",
                    master_info: SlideMasterInfo | None = None):
    """Render bullets with slide-type-specific, lower-noise layouts."""
    items = [item.strip() for item in content.items if item and item.strip()]
    if not items:
        return

    font_size = Pt(content.font_size) if content.font_size else config.FONT_BODY

    if slide_type == "agenda":
        _render_agenda_bullets(slide, pos, items, font_size, has_tpl, master_info)
        return

    if slide_type in ("executive_summary", "conclusion"):
        _render_summary_bullets(slide, pos, items, font_size, has_tpl, master_info)
        return

    _render_content_bullets(slide, pos, items, font_size, has_tpl, master_info)


def _render_agenda_bullets(slide, pos, items: list[str], font_size, has_tpl: bool,
                           master_info: SlideMasterInfo | None = None) -> None:
    cols = 2 if len(items) > 4 else 1
    rows = (len(items) + cols - 1) // cols
    gap_h = Emu(160000)
    gap_v = Emu(120000)
    card_w = (pos.width - gap_h * max(cols - 1, 0)) // max(cols, 1)
    card_h = min((pos.height - gap_v * max(rows - 1, 0)) // max(rows, 1), Emu(1000000))

    # Resolve a readable accent for the numbered badges so white-on-near-white
    # doesn't disappear on pastel templates (e.g. UAE Solar ACCENT_1 = EFF3E5).
    badge_hex = _pick_readable_accent_hex(master_info) if has_tpl else "4472C4"
    badge_rgb = _hex_to_rgb(badge_hex)

    for idx, item in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = pos.left + col * (card_w + gap_h)
        y = pos.top + row * (card_h + gap_v)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        _apply_light_surface_fill(card, has_tpl, brightness=0.98, fallback_hex="FBFCFE")
        add_shadow(card, preset="subtle")
        set_corner_radius(card, 8000)

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Emu(100000), y + Emu(120000), Emu(240000), Emu(240000))
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_rgb
        add_shadow(badge, preset="subtle")
        remove_outline(badge)
        _set_text_frame_text(
            badge.text_frame,
            str(idx + 1),
            font_size=Pt(11),
            bold=True,
            alignment=PP_ALIGN.CENTER,
            color_rgb=RGBColor(0xFF, 0xFF, 0xFF),
            font_name=config.FONT_NAME_PRIMARY,
        )

        tx = slide.shapes.add_textbox(x + Emu(420000), y + Emu(100000), card_w - Emu(540000), card_h - Emu(200000))
        # Use JetBrains Mono for code-like items, Inter for normal text
        item_font = config.FONT_NAME_MONO if _looks_like_code(item) else config.FONT_NAME_PRIMARY
        _populate_text_list(tx.text_frame, [item], font_size, font_name=item_font)


def _render_summary_bullets(slide, pos, items: list[str], font_size, has_tpl: bool,
                            master_info: SlideMasterInfo | None = None) -> None:
    """Render executive summary / conclusion bullets as rich cards with badges and connectors."""
    cols = 2 if len(items) >= 4 else 1
    rows = (len(items) + cols - 1) // cols
    gap_h = Emu(160000)
    gap_v = Emu(140000)
    card_w = (pos.width - gap_h * max(cols - 1, 0)) // max(cols, 1)
    card_h = min((pos.height - gap_v * max(rows - 1, 0)) // max(rows, 1), Emu(1600000))
    badge_sz = Emu(260000)
    accent_h = Emu(50000)

    # Readable accent for the top strip + numbered badge — avoids near-white
    # decorations on pastel templates.
    accent_hex = _pick_readable_accent_hex(master_info) if has_tpl else "4472C4"
    accent_rgb = _hex_to_rgb(accent_hex)

    for idx, item in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = pos.left + col * (card_w + gap_h)
        y = pos.top + row * (card_h + gap_v)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        _apply_light_surface_fill(card, has_tpl, brightness=0.97, fallback_hex="F8FAFD")
        add_shadow(card, preset="subtle")
        set_corner_radius(card, 8000)

        # Top accent bar (solid readable accent)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, accent_h)
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_rgb
        remove_outline(accent)

        # Numbered badge
        bx = x + Emu(120000)
        by = y - badge_sz // 3
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, badge_sz, badge_sz)
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent_rgb
        add_shadow(badge, preset="subtle")
        remove_outline(badge)
        _set_text_frame_text(
            badge.text_frame, str(idx + 1),
            font_size=Pt(11), bold=True, alignment=PP_ALIGN.CENTER,
            color_rgb=RGBColor(0xFF, 0xFF, 0xFF),
            font_name=config.FONT_NAME_PRIMARY,
        )

        tx = slide.shapes.add_textbox(x + Emu(140000), y + accent_h + Emu(80000),
                                       card_w - Emu(280000), card_h - accent_h - Emu(160000))
        item_font = config.FONT_NAME_MONO if _looks_like_code(item) else config.FONT_NAME_PRIMARY
        _populate_text_list(tx.text_frame, [item], font_size, font_name=item_font)

        # Connector line below card (between rows)
        if row < rows - 1:
            try:
                sep_y = y + card_h + gap_v // 2
                sep = slide.shapes.add_connector(
                    1, x + Emu(80000), sep_y, x + card_w - Emu(80000), sep_y
                )
                if has_tpl:
                    sep.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                else:
                    sep.line.color.rgb = _hex_to_rgb("B0C4DE")
                sep.line.width = Pt(0.75)
            except Exception:
                pass

    # Vertical separator between columns
    if cols > 1:
        vx = pos.left + card_w + gap_h // 2
        try:
            vsep = slide.shapes.add_connector(
                1, vx, pos.top + Emu(40000), vx, pos.top + rows * card_h + (rows - 1) * gap_v - Emu(40000)
            )
            if has_tpl:
                vsep.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            else:
                vsep.line.color.rgb = _hex_to_rgb("B0C4DE")
            vsep.line.width = Pt(0.5)
        except Exception:
            pass


def _render_content_bullets(slide, pos, items: list[str], font_size, has_tpl: bool,
                            master_info: SlideMasterInfo | None = None) -> None:
    """Render content bullets as a professional card grid with decorative elements.

    Each bullet becomes its own card with a numbered accent badge, accent top
    bar, and separator line — matching the shape density and visual quality of
    reference PPTXs (~25 shapes per content slide).
    """
    n = len(items)
    if n == 0:
        return

    bullet_font = (config.FONT_NAME_MONO
                    if any(_looks_like_code(item) for item in items)
                    else config.FONT_NAME_PRIMARY)

    # ── Layout: decide grid dimensions ──
    if n <= 2:
        cols, rows = n, 1
    elif n <= 4:
        cols, rows = min(n, 2), (n + 1) // 2
    else:
        cols, rows = min(n, 3), (n + 2) // 3

    gap_h = Emu(140000)
    gap_v = Emu(120000)
    card_w = (pos.width - gap_h * max(cols - 1, 0)) // max(cols, 1)
    card_h = min(
        (pos.height - gap_v * max(rows - 1, 0)) // max(rows, 1),
        Emu(2800000),
    )

    badge_sz = Emu(240000)  # numbered circle badge
    accent_bar_h = Emu(50000)  # top accent bar height

    # Readable accent for the top strip + numbered badge
    accent_hex = _pick_readable_accent_hex(master_info) if has_tpl else "4472C4"
    accent_rgb = _hex_to_rgb(accent_hex)

    for idx, item in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = pos.left + col * (card_w + gap_h)
        y = pos.top + row * (card_h + gap_v)

        # ── Card background rectangle ──
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        _apply_light_surface_fill(card, has_tpl, brightness=0.97, fallback_hex="F8FAFD")
        add_shadow(card, preset="subtle")

        # ── Top accent bar (solid readable accent) ──
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, accent_bar_h)
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_rgb
        remove_outline(accent)

        # ── Numbered badge (circle overlapping top-left corner) ──
        badge_x = x + Emu(100000)
        badge_y = y - badge_sz // 3
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y, badge_sz, badge_sz)
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent_rgb
        add_shadow(badge, preset="subtle")
        remove_outline(badge)
        _set_text_frame_text(
            badge.text_frame,
            str(idx + 1),
            font_size=Pt(11),
            bold=True,
            alignment=PP_ALIGN.CENTER,
            color_rgb=RGBColor(0xFF, 0xFF, 0xFF),
            font_name=config.FONT_NAME_PRIMARY,
        )

        # ── Text content inside card ──
        tx = slide.shapes.add_textbox(
            x + Emu(100000), y + accent_bar_h + Emu(60000),
            card_w - Emu(200000), card_h - accent_bar_h - Emu(120000),
        )
        _populate_text_list(tx.text_frame, [item], font_size, font_name=bullet_font)

        # ── Bottom separator line ──
        try:
            sep_y = y + card_h + gap_v // 2
            if row < rows - 1:  # horizontal between rows
                sep = slide.shapes.add_connector(
                    1, x + Emu(60000), sep_y, x + card_w - Emu(60000), sep_y
                )
                if has_tpl:
                    sep.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                    sep.line.width = Pt(0.75)
                else:
                    sep.line.color.rgb = _hex_to_rgb("B0C4DE")
                    sep.line.width = Pt(0.75)
        except Exception:
            pass

    # ── Vertical separators between columns ──
    for c in range(1, cols):
        vx = pos.left + c * (card_w + gap_h) - gap_h // 2
        try:
            vsep = slide.shapes.add_connector(
                1, vx, pos.top + Emu(60000), vx, pos.top + card_h * rows + gap_v * (rows - 1) - Emu(60000)
            )
            if has_tpl:
                vsep.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                vsep.line.width = Pt(0.5)
            else:
                vsep.line.color.rgb = _hex_to_rgb("B0C4DE")
                vsep.line.width = Pt(0.5)
        except Exception:
            pass


# ── Chart rendering ─────────────────────────────────────────────────

def _render_chart(slide, pos, content: ChartContent,
                  master_info: SlideMasterInfo | None = None, has_tpl: bool = False):
    """Render a native PowerPoint chart."""
    # Shadow container behind chart (charts can't have shadows directly)
    render_chart_container(slide, pos.left, pos.top, pos.width, pos.height, has_tpl)

    chart_type = CHART_TYPE_MAP.get(content.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    if content.chart_type == "scatter":
        chart_data = XyChartData()
        for series in content.series:
            s = chart_data.add_series(series.name)
            for i, val in enumerate(series.values):
                x = float(i)
                s.add_data_point(x, val)
    else:
        chart_data = CategoryChartData()
        chart_data.categories = content.categories
        for series in content.series:
            chart_data.add_series(series.name, tuple(series.values))

    graphic_frame = slide.shapes.add_chart(
        chart_type, pos.left, pos.top, pos.width, pos.height, chart_data
    )
    chart = graphic_frame.chart

    # Title
    if content.title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = content.title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.chart_title.text_frame.paragraphs[0].font.name = config.FONT_NAME_PRIMARY
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    # Legend — always show for multi-series, and for pie/doughnut
    if len(content.series) > 1 or content.chart_type in ("pie", "doughnut"):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        try:
            chart.legend.font.size = Pt(9)
            chart.legend.font.name = config.FONT_NAME_PRIMARY
        except Exception:
            pass

    # Pick a compact number format when the series values are large enough
    # that raw digits (e.g. 1,200,000,000) would crowd the chart. The custom
    # format below reads as "use T/B/M/K suffixes automatically".
    def _compact_num_format(series_list) -> str | None:
        try:
            vals = [abs(v) for s in series_list for v in s.values
                    if v is not None and not math.isnan(v)]
        except Exception:
            vals = []
        if not vals:
            return None
        peak = max(vals)
        if peak >= 1_000_000_000_000:
            return '[>=1000000000]0.0,,,"B";[>=1000000]0.0,,"M";[>=1000]0,"K";#,##0'
        if peak >= 1_000_000_000:
            return '[>=1000000000]0.0,,,"B";[>=1000000]0.0,,"M";[>=1000]0,"K";#,##0'
        if peak >= 1_000_000:
            return '[>=1000000]0.0,,"M";[>=1000]0,"K";#,##0'
        if peak >= 10_000:
            return '[>=1000]0,"K";#,##0'
        return None

    compact_fmt = _compact_num_format(content.series)

    # Data labels only when they help readability rather than creating clutter.
    try:
        plot = chart.plots[0]
        show_data_labels = (
            content.chart_type in ("pie", "doughnut")
            or (content.chart_type in ("bar", "column")
                and len(content.series) <= 2
                and len(content.categories) <= 10)
            or (len(content.series) == 1 and len(content.categories) <= 8)
        )
        plot.has_data_labels = show_data_labels
        if show_data_labels:
            data_labels = plot.data_labels
            data_labels.font.size = Pt(9)
            data_labels.font.name = config.FONT_NAME_MONO
            if content.chart_type in ("pie", "doughnut"):
                data_labels.number_format = '0%'
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            elif content.chart_type in ("bar", "column"):
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                if compact_fmt:
                    data_labels.number_format = compact_fmt
            elif content.chart_type == "line":
                data_labels.position = XL_LABEL_POSITION.ABOVE
                if compact_fmt:
                    data_labels.number_format = compact_fmt
    except Exception:
        pass  # some chart types don't support all label positions

    # Axis formatting (skip for pie/doughnut which have no axes)
    if content.chart_type not in ("pie", "doughnut"):
        try:
            # Category axis
            cat_axis = chart.category_axis
            cat_axis.has_major_gridlines = False
            cat_axis.tick_labels.font.size = Pt(9)
            cat_axis.tick_labels.font.name = config.FONT_NAME_PRIMARY
            # Value axis
            val_axis = chart.value_axis
            val_axis.has_major_gridlines = True
            val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            val_axis.has_minor_gridlines = False
            val_axis.tick_labels.font.size = Pt(9)
            val_axis.tick_labels.font.name = config.FONT_NAME_PRIMARY
            # Apply the same compact number format to axis ticks so giant
            # raw values don't monopolise the axis gutter.
            if compact_fmt:
                try:
                    val_axis.tick_labels.number_format = compact_fmt
                    val_axis.tick_labels.number_format_is_linked = False
                except Exception:
                    pass
            # Log scale when values span >2 orders of magnitude
            if getattr(content, "log_scale", False):
                try:
                    val_axis_el = val_axis._element
                    scaling = val_axis_el.find(qn("c:scaling"))
                    if scaling is not None:
                        existing = scaling.find(qn("c:logBase"))
                        if existing is None:
                            log_base = etree.SubElement(scaling, qn("c:logBase"))
                            log_base.set("val", "10")
                            # logBase must come before orientation per schema
                            orient = scaling.find(qn("c:orientation"))
                            if orient is not None:
                                scaling.remove(log_base)
                                scaling.insert(list(scaling).index(orient), log_base)
                except Exception as e:
                    _log.debug(f"Failed to set log scale: {e}")
        except Exception:
            pass  # scatter charts may have different axis structure

    # Color series using theme colors (full accent cycling for charts)
    for i, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        if has_tpl:
            fill.fore_color.theme_color = _ACCENT_THEME_COLORS[i % len(_ACCENT_THEME_COLORS)]
        else:
            fill.fore_color.rgb = _hex_to_rgb(_FALLBACK_ACCENT_HEX[i % len(_FALLBACK_ACCENT_HEX)])


# ── Table rendering ─────────────────────────────────────────────────

def _render_table(slide, pos, content: TableContent,
                  master_info: SlideMasterInfo | None = None,
                  has_tpl: bool = False):
    """Render a formatted table.

    Clamps table height so the last row never extends below the footer bar
    (``_sh - 900k EMU`` safe area) and, if necessary, drops excess rows.
    """
    cols = len(content.headers)
    if cols == 0:
        return

    # Reserve 900k EMU at the bottom for the footer/slide-number furniture.
    _FOOTER_RESERVE = 900000
    max_bottom = _sh - _FOOTER_RESERVE
    available_h = max(max_bottom - pos.top, 800000)  # min 0.88 in
    clamped_h = min(pos.height, available_h)

    # Estimate row height: header row ~55k EMU vertical padding + 11pt text,
    # data rows ~36k + 10pt text. Use conservative 420000 EMU per row as default.
    avg_row_h = 420000
    max_rows_fit = max(1, clamped_h // avg_row_h)

    # Header + data row cap
    data_rows = content.rows
    if len(data_rows) + 1 > max_rows_fit:
        _log.info(f"Trimming table from {len(data_rows)} data rows to {max_rows_fit - 1} to fit")
        data_rows = data_rows[: max(0, max_rows_fit - 1)]

    rows = len(data_rows) + 1  # +1 for header
    table_shape = slide.shapes.add_table(rows, cols, pos.left, pos.top, pos.width, clamped_h)
    table = table_shape.table

    # Re-bind data rows for the rest of the function
    content = TableContent(headers=content.headers, rows=data_rows, col_widths=content.col_widths)

    # Calculate column widths
    if content.col_widths:
        total = sum(content.col_widths)
        for i, w in enumerate(content.col_widths):
            if i < cols:
                table.columns[i].width = int(pos.width * w / total)
    else:
        col_width = pos.width // cols
        for i in range(cols):
            table.columns[i].width = col_width

    # Resolve header color from theme
    header_hex = _get_accent_hex(master_info, 0) if master_info else "1F4E79"
    # Derive a lighter shade for alternating rows
    alt_row_hex = _get_accent_hex(master_info, 0) if master_info else "D6E4F0"

    # Header row — pick contrast-aware text color (guard against UAE's pale accent)
    accent_hex = _get_accent_hex(master_info, 0) if master_info else header_hex
    header_text_hex = pick_text_color(accent_hex, large_text=True)
    header_text_rgb = _hex_to_rgb(header_text_hex)

    for col_idx, header in enumerate(content.headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(11)
        para.font.name = config.FONT_NAME_PRIMARY
        para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
        cell.fill.solid()
        # Use hex fill (derived from theme) so we can reliably pick text contrast.
        # Theme-color fill with auto-text can flip to light-on-light on pale accents.
        cell.fill.fore_color.rgb = _hex_to_rgb(accent_hex)
        para.font.color.rgb = header_text_rgb
        # Cell padding - generous for readability
        cell.margin_left = Emu(100000)
        cell.margin_right = Emu(100000)
        cell.margin_top = Emu(55000)
        cell.margin_bottom = Emu(55000)

    # Estimate max chars per cell based on column width — prevents mid-word wraps
    max_chars_per_cell: list[int] = []
    for i in range(cols):
        col_w_emu = table.columns[i].width or (pos.width // cols)
        col_w_in = col_w_emu / 914400
        # ~10 chars per inch at 10pt
        max_chars_per_cell.append(max(int(col_w_in * 11), 20))

    # Data rows
    for row_idx, row_data in enumerate(content.rows):
        for col_idx, value in enumerate(row_data):
            if col_idx < cols:
                cell = table.cell(row_idx + 1, col_idx)
                text = str(value).strip()
                # Shorten long numbers ($865,000,000 → $865M) before any
                # truncation so the abbreviation always fits.
                text = abbreviate_number(text)
                # Word-boundary truncate if too long for the column
                max_ch = max_chars_per_cell[col_idx] * 3  # allow up to 3 lines of wrap
                if len(text) > max_ch:
                    cut = text[:max_ch]
                    last_space = cut.rfind(' ')
                    if last_space >= max_ch * 0.6:
                        cut = cut[:last_space]
                    text = cut.rstrip(" ,.;:-") + "…"
                cell.text = text
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.name = config.FONT_NAME_PRIMARY
                para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                # Ultra-light alternating row shading (0.92 brightness)
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    if has_tpl:
                        cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
                        cell.fill.fore_color.brightness = 0.92
                    else:
                        cell.fill.fore_color.rgb = _hex_to_rgb("EDF2F9")
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                # Cell padding - compact for data rows
                cell.margin_left = Emu(100000)
                cell.margin_right = Emu(100000)
                cell.margin_top = Emu(36000)
                cell.margin_bottom = Emu(36000)
                # Predictable text layout
                _set_autofit(cell.text_frame)


# ── Shape rendering ─────────────────────────────────────────────────

def _render_shape(slide, pos, content: ShapeContent):
    """Render an auto shape."""
    shape_type = SHAPE_MAP.get(content.shape_type, MSO_SHAPE.ROUNDED_RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, pos.left, pos.top, pos.width, pos.height)

    if content.text:
        shape.text = content.text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        for para in tf.paragraphs:
            para.font.size = Pt(content.font_size) if content.font_size else config.FONT_BODY
            para.font.name = config.FONT_NAME_PRIMARY
            if content.bold:
                para.font.bold = True
            para.alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if content.fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(content.fill_color)
    else:
        shape.fill.background()  # transparent

    if content.line_color:
        shape.line.color.rgb = _hex_to_rgb(content.line_color)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()  # no border


# ── Infographic rendering ───────────────────────────────────────────

def _render_infographic(slide, pos, content: InfographicContent,
                       master_info: SlideMasterInfo | None = None,
                       has_tpl: bool = False):
    """Render an infographic based on type."""
    if content.infographic_type == "process_flow":
        _render_process_flow(slide, pos, content.items, master_info, has_tpl)
    elif content.infographic_type == "timeline":
        _render_timeline(slide, pos, content.items, has_tpl)
    elif content.infographic_type == "comparison":
        _render_comparison(slide, pos, content.items, master_info, has_tpl)
    elif content.infographic_type == "kpi_cards":
        _render_kpi_cards(slide, pos, content.items, master_info, has_tpl)
    elif content.infographic_type == "hierarchy":
        _render_hierarchy(slide, pos, content.items, has_tpl)
    elif content.infographic_type == "icon_list":
        _render_icon_list(slide, pos, content.items, master_info, has_tpl)
    elif content.infographic_type == "stat_grid":
        _render_stat_grid(slide, pos, content.items, master_info, has_tpl)
    elif content.infographic_type == "hero_number":
        _render_hero_number(slide, pos, content.items, master_info, has_tpl)
    elif content.infographic_type == "pull_quote":
        _render_pull_quote(slide, pos, content.items, master_info, has_tpl)


def _render_process_flow(slide, pos, items, master_info: SlideMasterInfo | None = None,
                         has_tpl: bool = False):
    """Render a process flow with rounded rectangles, step number overlays, and arrow connectors.

    Visual discipline:
    - Single brand accent (ACCENT_1) — no rainbow cycling.
    - Cards use the accent at full saturation so step number legibility is
      guaranteed; text color is contrast-picked against the effective fill.
    - Shadow + corner radius for a modern card feel.
    """
    items = items[:config.MAX_PROCESS_ITEMS]
    n = len(items)
    if n == 0:
        return

    # Resolve the single accent hex that drives every card
    accent_hex = (
        master_info.theme_colors.accents()[0]
        if master_info else FLOW_ACCENT_HEX[0]
    )

    arrow_gap = Emu(260000)
    usable_w = pos.width - arrow_gap * max(n - 1, 0)
    item_width = min(usable_w // max(n, 1), Emu(3200000))

    total_w = item_width * n + arrow_gap * max(n - 1, 0)
    x_offset = pos.left + (pos.width - total_w) // 2

    item_height = Emu(min(pos.height, 1400000))
    y_center = pos.top + (pos.height - item_height) // 2
    step_circle_size = Emu(260000)

    for i, item in enumerate(items):
        x = x_offset + i * (item_width + arrow_gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y_center, item_width, item_height
        )
        if has_tpl:
            style_card(shape, theme_color=MSO_THEME_COLOR.ACCENT_1,
                       shadow_preset="card", corner_radius=8000)
        else:
            c1 = _hex_to_rgb(accent_hex)
            c2 = RGBColor(max(c1[0] - 20, 0), max(c1[1] - 20, 0), max(c1[2] - 20, 0))
            style_card(shape, gradient_stops=[(0.0, c1), (1.0, c2)],
                       shadow_preset="card", corner_radius=8000)

        # Step number badge (white circle overlay)
        sc_x = x - step_circle_size // 3
        sc_y = y_center - step_circle_size // 3
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, sc_x, sc_y, step_circle_size, step_circle_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.word_wrap = False
        cp = ctf.paragraphs[0]
        cp.text = str(i + 1)
        cp.font.size = Pt(12)
        cp.font.name = config.FONT_NAME_MONO
        cp.font.bold = True
        cp.font.color.rgb = _hex_to_rgb(accent_hex)
        cp.alignment = PP_ALIGN.CENTER

        # Contrast-aware text color based on the ACTUAL rendered fill
        txt_hex = pick_text_color(accent_hex, large_text=True)
        txt_rgb = _hex_to_rgb(txt_hex)
        sub_rgb = _hex_to_rgb(darken_hex(txt_hex, 0.08))

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = config.TF_MARGIN_LEFT
        tf.margin_right = config.TF_MARGIN_RIGHT
        tf.margin_top = Emu(100000)
        tf.margin_bottom = Emu(80000)
        _set_autofit(tf, shrink_ok=True)

        p = tf.paragraphs[0]
        p.text = item.title[:50]
        p.font.size = Pt(12)
        p.font.name = config.FONT_NAME_PRIMARY
        p.font.bold = True
        p.font.color.rgb = txt_rgb
        p.alignment = PP_ALIGN.CENTER

        if item.description:
            p2 = tf.add_paragraph()
            p2.text = item.description[:config.MAX_INFOGRAPHIC_DESC]
            p2.font.size = Pt(10)
            p2.font.name = config.FONT_NAME_PRIMARY
            p2.font.color.rgb = sub_rgb
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(8)

        # ── Arrow connector (real connector shape + chevron indicator) ──
        if i < n - 1:
            ax_start = x + item_width
            ax_end = ax_start + arrow_gap
            ay = y_center + item_height // 2
            try:
                conn = slide.shapes.add_connector(
                    1, ax_start, ay, ax_end, ay
                )
                if has_tpl:
                    conn.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                else:
                    conn.line.color.rgb = _hex_to_rgb("4472C4")
                conn.line.width = Pt(2)
            except Exception:
                pass
            # Small chevron at midpoint of arrow
            chev_sz = Emu(160000)
            chev_x = ax_start + (arrow_gap - chev_sz) // 2
            chev_y = ay - chev_sz // 2
            chev = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON, chev_x, chev_y, chev_sz, chev_sz
            )
            if has_tpl:
                chev.fill.solid()
                chev.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
            else:
                chev.fill.solid()
                chev.fill.fore_color.rgb = _hex_to_rgb("4472C4")
            chev.line.fill.background()


def _render_timeline(slide, pos, items, has_tpl: bool = False):
    """Render a horizontal timeline with alternating above/below labels.

    Layout discipline:
    - Max 5 items (more creates unreadable crowding; caller should split).
    - Value (year/date) shown INSIDE the circle only — never duplicated above.
    - Labels placed with ~900k EMU vertical clearance from the axis line.
    - Title + description drawn in a single "label block" to avoid overlap.
    """
    items = list(items)[:5]
    n = len(items)
    if n == 0:
        return

    # Horizontal axis line at vertical center of pos
    line_y = pos.top + pos.height // 2

    # Nodes
    node_gap = pos.width // max(n, 1)
    circle_size = Emu(340000)

    # Reserve vertical space for labels — need enough headroom above/below the line.
    # Each label block = title (≈300k) + description (≈480k) + gap (150k) = ~930k EMU.
    label_block_h = Emu(900000)
    clearance = Emu(80000)

    connector = slide.shapes.add_connector(
        1, pos.left + Emu(100000), line_y,
        pos.left + pos.width - Emu(100000), line_y,
    )
    if has_tpl:
        connector.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    else:
        connector.line.color.rgb = _hex_to_rgb("2E75B6")
    connector.line.width = Pt(2.5)

    for i, item in enumerate(items):
        cx = pos.left + i * node_gap + node_gap // 2

        # Circle marker
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx - circle_size // 2, line_y - circle_size // 2,
            circle_size, circle_size
        )
        if has_tpl:
            style_numbered_circle(circle, MSO_THEME_COLOR.ACCENT_1)
        else:
            circle.fill.solid()
            circle.fill.fore_color.rgb = _hex_to_rgb("4472C4")
            add_shadow(circle, preset="subtle")
            remove_outline(circle)

        # Circle label: prefer value (year/date) else step number
        circle_label = str(i + 1)
        value_str = (item.value or "").strip() if item.value else ""
        if value_str and len(value_str) <= 8:
            circle_label = value_str
        tf = circle.text_frame
        tf.word_wrap = False
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.text = circle_label
        p.font.size = Pt(9 if len(circle_label) > 4 else 11)
        p.font.name = config.FONT_NAME_MONO
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # Alternate label placement (above / below axis)
        is_above = (i % 2 == 0)
        if is_above:
            label_top = line_y - circle_size // 2 - clearance - label_block_h
        else:
            label_top = line_y + circle_size // 2 + clearance

        # Connecting stub: short dashed line from circle edge to label block
        stub_start = line_y - circle_size // 2 if is_above else line_y + circle_size // 2
        stub_end = label_top + label_block_h if is_above else label_top
        try:
            seg = slide.shapes.add_connector(1, cx, stub_start, cx, stub_end)
            if has_tpl:
                seg.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            else:
                seg.line.color.rgb = _hex_to_rgb("4472C4")
            seg.line.width = Pt(0.75)
            seg.line.dash_style = 2
        except Exception:
            pass

        # Single label block — title + description inside one textbox
        label_w = min(node_gap - Emu(80000), Emu(2400000))
        label_x = cx - label_w // 2
        title_text = (item.title or "").strip()[:60]
        desc_text = (item.description or "").strip()[:config.MAX_INFOGRAPHIC_DESC]

        # If value is long (date string), put it above title instead of in circle
        if value_str and len(value_str) > 8:
            composite = [value_str, title_text]
        else:
            composite = [title_text]
        if desc_text:
            composite.append(desc_text)

        tb = slide.shapes.add_textbox(label_x, label_top, label_w, label_block_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(40000)
        tf.margin_right = Emu(40000)
        tf.margin_top = Emu(20000)
        tf.margin_bottom = Emu(20000)
        tf.auto_size = MSO_AUTO_SIZE.NONE

        for j, txt in enumerate(composite):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = txt
            para.alignment = PP_ALIGN.CENTER
            para.font.name = config.FONT_NAME_PRIMARY
            if j == 0 and value_str and len(value_str) > 8:
                para.font.size = Pt(10)
                para.font.bold = True
                if has_tpl:
                    para.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                else:
                    para.font.color.rgb = _hex_to_rgb("2B5797")
            elif (j == 0 and not (value_str and len(value_str) > 8)) or (j == 1 and value_str and len(value_str) > 8):
                # This is the title paragraph
                para.font.size = Pt(11)
                para.font.bold = True
            else:
                para.font.size = Pt(9)
                para.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
                para.space_before = Pt(3)


def _render_comparison(slide, pos, items, master_info: SlideMasterInfo | None = None,
                       has_tpl: bool = False):
    """Render side-by-side comparison cards with a single brand accent.

    Uses one accent color (ACCENT_1) across every card — the Common Mistakes
    guide explicitly flagged rainbow cycling as "random colors". Visual
    variety comes from shadow + subtle gradient, not from colour changes.
    """
    n = len(items)
    if n == 0:
        return
    accent_hex = (
        master_info.theme_colors.accents()[0]
        if master_info else CMP_ACCENT_HEX[0]
    )

    gap = config.SHAPE_GAP
    card_width = (pos.width - gap * max(n - 1, 0)) // max(n, 1)
    card_height = pos.height

    txt_hex = pick_text_color(accent_hex, large_text=True)
    txt_rgb = _hex_to_rgb(txt_hex)
    sub_rgb = _hex_to_rgb(darken_hex(txt_hex, 0.08))

    for i, item in enumerate(items):
        x = pos.left + i * (card_width + gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, pos.top, card_width, card_height
        )
        if has_tpl:
            style_card(card, theme_color=MSO_THEME_COLOR.ACCENT_1,
                       shadow_preset="card", corner_radius=8000)
        else:
            c1 = _hex_to_rgb(accent_hex)
            c2 = RGBColor(max(c1[0] - 20, 0), max(c1[1] - 20, 0), max(c1[2] - 20, 0))
            style_card(card, gradient_stops=[(0.0, c1), (1.0, c2)],
                       shadow_preset="card", corner_radius=8000)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = config.TF_MARGIN_LEFT
        tf.margin_right = config.TF_MARGIN_RIGHT
        tf.margin_top = Emu(100000)
        tf.margin_bottom = Emu(80000)
        _set_autofit(tf, shrink_ok=True)

        p = tf.paragraphs[0]
        p.text = item.title[:60]
        p.font.size = Pt(13)
        p.font.name = config.FONT_NAME_PRIMARY
        p.font.bold = True
        p.font.color.rgb = txt_rgb
        p.alignment = PP_ALIGN.CENTER

        if item.description:
            p2 = tf.add_paragraph()
            p2.text = item.description[:config.MAX_CMP_DESC]
            p2.font.size = Pt(10)
            p2.font.name = config.FONT_NAME_PRIMARY
            p2.font.color.rgb = sub_rgb
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(10)

        if item.value:
            p3 = tf.add_paragraph()
            p3.text = item.value
            p3.font.size = Pt(20)
            p3.font.name = config.FONT_NAME_MONO
            p3.font.bold = True
            p3.font.color.rgb = txt_rgb
            p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(14)


def _render_kpi_cards(slide, pos, items, master_info: SlideMasterInfo | None = None,
                      has_tpl: bool = False):
    """Render KPI metric cards with value, title, and vertical centering.

    Single-accent brand discipline — every card uses ACCENT_1 at full
    saturation. Variety comes from the metric values themselves, not from
    different colours. Text color is contrast-picked against the actual
    rendered fill so pale-accent templates (like UAE's #EFF3E5) still
    produce readable cards.
    """
    items = items[:config.MAX_KPI_ITEMS]
    n = len(items)
    if n == 0:
        return

    accent_hex = (
        master_info.theme_colors.accents()[0]
        if master_info else KPI_ACCENT_HEX[0]
    )

    gap = config.SHAPE_GAP
    card_width = (pos.width - gap * max(n - 1, 0)) // max(n, 1)
    card_height = min(pos.height, Emu(2200000))
    y = pos.top + (pos.height - card_height) // 2
    stripe_h = Emu(60000)

    # One contrast-picked text color used for every card
    ktxt = _hex_to_rgb(pick_text_color(accent_hex, large_text=True))
    ksub = _hex_to_rgb(darken_hex(pick_text_color(accent_hex), 0.08))

    for i, item in enumerate(items):
        x = pos.left + i * (card_width + gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height
        )
        if has_tpl:
            style_card(card, theme_color=MSO_THEME_COLOR.ACCENT_1,
                       shadow_preset="card", corner_radius=8000)
        else:
            c1 = _hex_to_rgb(accent_hex)
            c2 = RGBColor(max(c1[0] - 25, 0), max(c1[1] - 25, 0), max(c1[2] - 25, 0))
            style_card(card, gradient_stops=[(0.0, c1), (1.0, c2)],
                       shadow_preset="card", corner_radius=8000)

        # Top frosted stripe for depth
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, card_width, stripe_h
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        sp_el = stripe._element.spPr
        if sp_el is not None:
            fill_el = sp_el.find(qn("a:solidFill"))
            if fill_el is not None:
                srgb = fill_el.find(qn("a:srgbClr"))
                if srgb is not None:
                    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
                    alpha_el.set("val", "30000")
        stripe.line.fill.background()

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Emu(stripe_h + 80000)
        tf.margin_left = config.TF_MARGIN_LEFT
        tf.margin_right = config.TF_MARGIN_RIGHT
        tf.margin_bottom = Emu(80000)
        _set_autofit(tf, shrink_ok=True)

        _ktxt = ktxt
        _ksub = ksub

        # Big value — prominent hero metric (40pt+ for visual hierarchy)
        p = tf.paragraphs[0]
        p.text = item.value or ""
        p.font.size = Pt(42)
        p.font.name = config.FONT_NAME_MONO
        p.font.bold = True
        p.font.color.rgb = _ktxt
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

        # ── Real separator line between value and label ──
        try:
            sep_line = slide.shapes.add_connector(
                1, x + Emu(200000), y + card_height // 2 + Emu(60000),
                x + card_width - Emu(200000), y + card_height // 2 + Emu(60000)
            )
            if has_tpl:
                sep_line.line.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
            else:
                sep_line.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            sep_line.line.width = Pt(1)
        except Exception:
            pass

        # Text-based fallback separator
        sep = tf.add_paragraph()
        sep.text = ""
        sep.font.size = Pt(4)
        sep.alignment = PP_ALIGN.CENTER
        sep.space_before = Pt(4)

        # Label below separator
        p2 = tf.add_paragraph()
        p2.text = item.title
        p2.font.size = Pt(12)
        p2.font.name = config.FONT_NAME_PRIMARY
        p2.font.color.rgb = _ksub
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)

        # Description if present
        if item.description:
            p3 = tf.add_paragraph()
            p3.text = item.description[:config.MAX_INFOGRAPHIC_DESC]
            p3.font.size = Pt(9)
            p3.font.name = config.FONT_NAME_PRIMARY
            p3.font.color.rgb = _ksub
            p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(4)


def _render_hierarchy(slide, pos, items, has_tpl: bool = False):
    """Render a hierarchy with indented boxes and connectors."""
    n = len(items)
    if n == 0:
        return
    gap = Emu(80000)
    item_height = min((pos.height - gap * max(n - 1, 0)) // max(n, 1), Emu(900000))

    for i, item in enumerate(items):
        indent = Emu(i * 220000)
        y = pos.top + i * (item_height + gap)
        w = pos.width - indent

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, pos.left + indent, y, w, item_height
        )
        _apply_accent_fill(shape, i, has_tpl)
        add_shadow(shape, preset="subtle")
        set_corner_radius(shape, 8000)
        remove_outline(shape)

        # Vertical connector to next level
        if i < n - 1:
            cx = pos.left + indent + Emu(120000)
            cy = y + item_height
            next_y = pos.top + (i + 1) * (item_height + gap)
            connector = slide.shapes.add_connector(
                1, cx, cy, cx, next_y
            )
            connector.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
            connector.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = config.TF_MARGIN_LEFT
        tf.margin_right = config.TF_MARGIN_RIGHT
        tf.margin_top = Emu(60000)
        tf.margin_bottom = Emu(40000)
        _set_autofit(tf, shrink_ok=True)

        # Contrast-aware text color
        _hfb = _FALLBACK_ACCENT_HEX[i % len(_FALLBACK_ACCENT_HEX)]
        _htxt = _hex_to_rgb(pick_text_color(_hfb))
        _hsub = _hex_to_rgb(darken_hex(pick_text_color(_hfb), 0.08))

        p = tf.paragraphs[0]
        p.text = item.title
        p.font.size = Pt(12)
        p.font.name = config.FONT_NAME_PRIMARY
        p.font.bold = True
        p.font.color.rgb = _htxt
        p.alignment = PP_ALIGN.LEFT

        if item.description:
            p2 = tf.add_paragraph()
            p2.text = item.description[:config.MAX_INFOGRAPHIC_DESC]
            p2.font.size = Pt(10)
            p2.font.name = config.FONT_NAME_PRIMARY
            p2.font.color.rgb = _hsub
            p2.space_before = Pt(4)


# ── New archetypes introduced by the visual-overhaul iteration ──────

def _truncate_at_word(s: str, max_chars: int) -> str:
    """Cut *s* at a word boundary near *max_chars*, appending '…' if cut."""
    s = (s or "").strip()
    if len(s) <= max_chars:
        return s
    cut = s[:max_chars]
    idx = cut.rfind(" ")
    if idx >= max_chars * 0.55:
        cut = cut[:idx]
    return cut.rstrip(" ,.;:-") + "…"


def _render_icon_list(slide, pos, items, master_info: SlideMasterInfo | None = None,
                       has_tpl: bool = False):
    """Icon + title + description rows stacked vertically.

    Replaces the monotonous ``numbered-card grid`` that every content slide
    used to default to. Each row gets an auto-picked icon, a bold title and
    a short description — far more editorial and scannable.
    """
    from .icons import icon_for_keyword, draw_icon

    items = items[:6]
    n = len(items)
    if n == 0:
        return

    gap = Emu(110000)
    row_h = (pos.height - gap * max(n - 1, 0)) // max(n, 1)
    # Raise the per-row ceiling so titles + descriptions don't get clipped —
    # the old ~1.09" cap meant 2-line titles overran the visible area.
    row_h = min(row_h, Emu(1300000))

    icon_size = min(row_h - Emu(200000), Emu(700000))
    icon_col_w = icon_size + Emu(260000)
    text_x = pos.left + icon_col_w
    text_w = pos.width - icon_col_w - Emu(60000)

    # Shrink font slightly when rows are short (n=5-6) so wrapping stays inside the cell
    title_pt = 13 if n <= 4 else 12
    desc_pt = 10 if n <= 4 else 9

    for i, item in enumerate(items):
        y = pos.top + i * (row_h + gap)

        # Subtle underline separator (except last row)
        if i < n - 1:
            try:
                sep_y = y + row_h + gap // 2
                sep = slide.shapes.add_connector(
                    1, pos.left, sep_y, pos.left + pos.width, sep_y
                )
                if has_tpl:
                    sep.line.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                else:
                    sep.line.color.rgb = RGBColor(0xD0, 0xD6, 0xE0)
                sep.line.width = Pt(0.5)
            except Exception:
                pass

        # Icon — prefer explicit field, else auto-pick from title
        icon_name = item.icon or icon_for_keyword(item.title)
        icon_y = y + (row_h - icon_size) // 2
        draw_icon(slide, icon_name, pos.left + Emu(20000), icon_y,
                  icon_size, master_info, has_tpl)

        # Text block — title (bold) + description (muted)
        tb = slide.shapes.add_textbox(text_x, y, text_w, row_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(20000)
        tf.margin_right = Emu(20000)
        tf.margin_top = Emu(30000)
        tf.margin_bottom = Emu(20000)
        tf.auto_size = MSO_AUTO_SIZE.NONE
        # Vertical-centre the text so the row looks balanced next to the icon
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_para = tf.paragraphs[0]
        title_para.text = _truncate_at_word((item.title or ""), 95)
        title_para.font.size = Pt(title_pt)
        title_para.font.bold = True
        title_para.font.name = config.FONT_NAME_PRIMARY
        title_para.alignment = PP_ALIGN.LEFT

        desc = (item.description or "").strip()
        if item.value and item.value.strip():
            # Append value after description as a bold inline tag
            vstr = item.value.strip()
            desc = f"{vstr} · {desc}" if desc else vstr
        if desc:
            dp = tf.add_paragraph()
            dp.text = _truncate_at_word(desc, 160)
            dp.font.size = Pt(desc_pt)
            dp.font.name = config.FONT_NAME_PRIMARY
            dp.font.color.rgb = RGBColor(0x5A, 0x5A, 0x5A)
            dp.alignment = PP_ALIGN.LEFT
            dp.space_before = Pt(3)


def _render_stat_grid(slide, pos, items, master_info: SlideMasterInfo | None = None,
                      has_tpl: bool = False):
    """3-4 hero metrics in a balanced grid.

    Differs from ``kpi_cards`` by emphasising raw numeric dominance: no
    coloured card background, just huge metric values + small labels
    arranged cleanly. Matches the Common-Mistakes guide note "Numbers are
    not visually dominant enough".
    """
    items = [it for it in items if it.value and it.value.strip()][:4]
    n = len(items)
    if n == 0:
        return

    # Use the contrast-aware accent so pastel ACCENT_1 (e.g. UAE `EFF3E5`)
    # doesn't render the hero number invisible on a white slide.
    accent_hex = _pick_readable_accent_hex(master_info) if has_tpl else FALLBACK_ACCENT_HEX[0]
    accent_rgb = _hex_to_rgb(accent_hex)

    gap = Emu(240000)
    cell_w = (pos.width - gap * max(n - 1, 0)) // max(n, 1)
    cell_h = min(pos.height, Emu(2200000))
    y = pos.top + (pos.height - cell_h) // 2

    for i, item in enumerate(items):
        x = pos.left + i * (cell_w + gap)

        # Thin top accent bar for polish (solid readable accent)
        bar_h = Emu(50000)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + cell_w // 4, y, cell_w // 2, bar_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_rgb
        remove_outline(bar)

        # Big numeric value — always abbreviated so "1,200,000,000" never
        # occupies a whole card (and drop the 12-char cap when the shortened
        # form fits more context like "$1.2B USD").
        val_text = abbreviate_number(item.value.strip())
        if len(val_text) > 14:
            val_text = val_text[:14]

        val_box = slide.shapes.add_textbox(
            x, y + bar_h + Emu(40000), cell_w, cell_h // 2
        )
        vtf = val_box.text_frame
        vtf.word_wrap = False
        vtf.margin_left = Emu(0)
        vtf.margin_right = Emu(0)
        vtf.margin_top = Emu(0)
        vtf.margin_bottom = Emu(0)
        vp = vtf.paragraphs[0]
        vp.text = val_text
        vp.font.size = Pt(48)
        vp.font.bold = True
        vp.font.name = config.FONT_NAME_PRIMARY
        vp.font.color.rgb = accent_rgb
        vp.alignment = PP_ALIGN.CENTER

        # Label (title) + optional description
        label_box = slide.shapes.add_textbox(
            x, y + cell_h // 2 + Emu(100000), cell_w, cell_h // 2 - Emu(120000)
        )
        ltf = label_box.text_frame
        ltf.word_wrap = True
        ltf.margin_left = Emu(40000)
        ltf.margin_right = Emu(40000)
        ltf.margin_top = Emu(10000)
        ltf.margin_bottom = Emu(20000)
        ltf.auto_size = MSO_AUTO_SIZE.NONE

        lp = ltf.paragraphs[0]
        lp.text = (item.title or "").strip()[:40]
        lp.font.size = Pt(12)
        lp.font.bold = True
        lp.font.name = config.FONT_NAME_PRIMARY
        lp.alignment = PP_ALIGN.CENTER
        lp.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

        if item.description:
            dp = ltf.add_paragraph()
            dp.text = item.description.strip()[:config.MAX_INFOGRAPHIC_DESC]
            dp.font.size = Pt(10)
            dp.font.name = config.FONT_NAME_PRIMARY
            dp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            dp.alignment = PP_ALIGN.CENTER
            dp.space_before = Pt(4)


def _render_hero_number(slide, pos, items, master_info: SlideMasterInfo | None = None,
                        has_tpl: bool = False):
    """One gigantic metric on the left + narrative sidebar on the right.

    Use when a section is anchored on a single dominant statistic.
    """
    if not items:
        return
    item = items[0]

    # Contrast-aware accent (same path used by footer/stat_grid) so the hero
    # number remains visible on pastel templates.
    accent_hex = _pick_readable_accent_hex(master_info) if has_tpl else FALLBACK_ACCENT_HEX[0]
    accent_rgb = _hex_to_rgb(accent_hex)
    # Keep a separate "raw first accent" for the decorative arc so it still
    # follows the template's hue even if it's pastel.
    arc_base_hex = (
        master_info.theme_colors.accents()[0]
        if master_info and master_info.theme_colors.accents() else accent_hex
    )

    # Left 45% = hero number, right 55% = narrative
    left_w = int(pos.width * 0.45)
    right_x = pos.left + left_w + Emu(220000)
    right_w = pos.width - left_w - Emu(220000)

    # Decorative accent arc behind the number (large oval, clipped visually)
    arc_d = min(left_w, pos.height) + Emu(400000)
    arc_x = pos.left - arc_d // 4
    arc_y = pos.top + (pos.height - arc_d) // 2
    arc = slide.shapes.add_shape(MSO_SHAPE.OVAL, arc_x, arc_y, arc_d, arc_d)
    arc.fill.solid()
    arc.fill.fore_color.rgb = _hex_to_rgb(lighten_hex_simple(arc_base_hex, 0.82))
    remove_outline(arc)

    # Hero value — abbreviate and fall back to title when value is missing.
    raw_val = (item.value or "").strip() or (item.title or "").strip()
    short_val = abbreviate_number(raw_val)
    if len(short_val) > 14:
        short_val = short_val[:14]

    val_box = slide.shapes.add_textbox(
        pos.left, pos.top + pos.height // 4, left_w, pos.height // 2
    )
    vtf = val_box.text_frame
    vtf.word_wrap = True
    vtf.margin_left = Emu(20000)
    vtf.margin_right = Emu(20000)
    vtf.margin_top = Emu(0)
    vtf.margin_bottom = Emu(0)
    vp = vtf.paragraphs[0]
    vp.text = short_val
    vp.font.size = Pt(84)
    vp.font.bold = True
    vp.font.name = config.FONT_NAME_PRIMARY
    vp.font.color.rgb = accent_rgb
    vp.alignment = PP_ALIGN.LEFT

    # Sub-label under the hero number
    label = item.title if item.value else ""
    if label:
        sub_box = slide.shapes.add_textbox(
            pos.left, pos.top + pos.height // 4 + Emu(900000),
            left_w, Emu(400000),
        )
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = label.strip()[:60]
        sp.font.size = Pt(14)
        sp.font.bold = True
        sp.font.name = config.FONT_NAME_PRIMARY
        sp.alignment = PP_ALIGN.LEFT
        sp.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Right-side narrative
    narr_box = slide.shapes.add_textbox(
        right_x, pos.top + Emu(100000), right_w, pos.height - Emu(200000)
    )
    ntf = narr_box.text_frame
    ntf.word_wrap = True
    ntf.margin_left = Emu(40000)
    ntf.margin_right = Emu(40000)
    ntf.margin_top = Emu(20000)
    ntf.margin_bottom = Emu(20000)

    # Paragraph 1: item title (if value was used as hero) or description
    desc_text = (item.description or "").strip()
    paragraphs: list[tuple[str, int, bool, RGBColor]] = []
    if item.value and item.title:
        paragraphs.append((item.title.strip()[:80], 16, True, RGBColor(0x22, 0x22, 0x22)))
    if desc_text:
        paragraphs.append((desc_text[:240], 12, False, RGBColor(0x55, 0x55, 0x55)))

    # Bonus narrative rows — additional items' title + description
    for extra in items[1:4]:
        if extra.title:
            prefix = ""
            if extra.value:
                prefix = f"{extra.value.strip()} · "
            paragraphs.append((
                f"{prefix}{extra.title.strip()[:70]}",
                11, False, RGBColor(0x44, 0x44, 0x44),
            ))

    for j, (txt, sz, bold, col) in enumerate(paragraphs):
        para = ntf.paragraphs[0] if j == 0 else ntf.add_paragraph()
        para.text = txt
        para.font.size = Pt(sz)
        para.font.bold = bold
        para.font.name = config.FONT_NAME_PRIMARY
        para.font.color.rgb = col
        para.alignment = PP_ALIGN.LEFT
        if j > 0:
            para.space_before = Pt(8)


def _render_pull_quote(slide, pos, items, master_info: SlideMasterInfo | None = None,
                       has_tpl: bool = False):
    """Editorial pull-quote: large italic quote with decorative marks + attribution."""
    if not items:
        return
    item = items[0]

    accent_hex = (
        master_info.theme_colors.accents()[0]
        if master_info else FALLBACK_ACCENT_HEX[0]
    )
    accent_rgb = _hex_to_rgb(accent_hex)

    # Open-quote mark (large decorative, top-left)
    q_size = Emu(900000)
    q_mark = slide.shapes.add_textbox(
        pos.left, pos.top - Emu(200000), q_size, q_size
    )
    qtf = q_mark.text_frame
    qtf.word_wrap = False
    qtf.margin_left = Emu(0)
    qtf.margin_top = Emu(0)
    qp = qtf.paragraphs[0]
    qp.text = "“"
    qp.font.size = Pt(120)
    qp.font.bold = True
    qp.font.name = config.FONT_NAME_PRIMARY
    if has_tpl:
        qp.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        qp.font.color.brightness = 0.4
    else:
        qp.font.color.rgb = _hex_to_rgb(lighten_hex_simple(accent_hex, 0.5))
    qp.alignment = PP_ALIGN.LEFT

    # Quote body — italic, large, centered
    body_x = pos.left + Emu(600000)
    body_y = pos.top + Emu(300000)
    body_w = pos.width - Emu(700000)
    body_h = int(pos.height * 0.6)
    body = slide.shapes.add_textbox(body_x, body_y, body_w, body_h)
    btf = body.text_frame
    btf.word_wrap = True
    btf.margin_left = Emu(40000)
    btf.margin_right = Emu(40000)
    bp = btf.paragraphs[0]
    # Combine title and description into the quote body
    quote_text = (item.title or "").strip()
    if item.description:
        quote_text = (quote_text + ". " + item.description.strip()).strip(". ")
    bp.text = quote_text[:220]
    bp.font.size = Pt(24)
    bp.font.italic = True
    bp.font.name = config.FONT_NAME_PRIMARY
    bp.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    bp.alignment = PP_ALIGN.LEFT

    # Attribution underline + text
    attr_y = pos.top + int(pos.height * 0.78)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        pos.left + Emu(100000), attr_y,
        Emu(600000), Emu(20000),
    )
    line.fill.solid()
    if has_tpl:
        line.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    else:
        line.fill.fore_color.rgb = accent_rgb
    remove_outline(line)

    if item.value:
        attr_box = slide.shapes.add_textbox(
            pos.left + Emu(100000), attr_y + Emu(60000),
            pos.width - Emu(200000), Emu(400000),
        )
        atf = attr_box.text_frame
        ap = atf.paragraphs[0]
        ap.text = f"— {item.value.strip()}"
        ap.font.size = Pt(12)
        ap.font.bold = True
        ap.font.name = config.FONT_NAME_PRIMARY
        ap.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def lighten_hex_simple(hex_color: str, amount: float) -> str:
    """Module-local hex lightener (avoids circular import with color_utils)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * amount)
    g = int(g + (255 - g) * amount)
    b = int(b + (255 - b) * amount)
    return f"{r:02X}{g:02X}{b:02X}"
