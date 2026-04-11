"""Reusable visual components for professional PPTX slides.

Each component creates a self-contained visual element on a slide
(e.g., KPI card, content card, accent bar, numbered bullet row).
Uses the ``drawingml_effects`` module for shadows, gradients, and
corner radius.

All components respect ``has_tpl`` to decide between theme colours
(auto-adapts to any Slide Master) and fallback hex values.
"""

from __future__ import annotations

import logging
from typing import Sequence

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Emu, Pt

from . import config
from .drawingml_effects import (
    add_shadow,
    add_gradient,
    add_theme_gradient,
    remove_outline,
    set_corner_radius,
    style_card,
    style_accent_bar,
    style_numbered_circle,
)

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Theme colour helpers (mirrors pptx_renderer's constants)
# ---------------------------------------------------------------------------
_ACCENT_THEME = [
    MSO_THEME_COLOR.ACCENT_1, MSO_THEME_COLOR.ACCENT_2,
    MSO_THEME_COLOR.ACCENT_3, MSO_THEME_COLOR.ACCENT_4,
    MSO_THEME_COLOR.ACCENT_5, MSO_THEME_COLOR.ACCENT_6,
]

_FALLBACK_HEX = [
    "4472C4", "ED7D31", "A5A5A5", "FFC000", "5B9BD5", "70AD47",
]


def _hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_autofit(tf) -> None:
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.word_wrap = True


# ---------------------------------------------------------------------------
# KPI Card  —  large hero number + label + optional description
# ---------------------------------------------------------------------------

def render_kpi_card(
    slide,
    x: int, y: int, w: int, h: int,
    value: str,
    label: str,
    description: str = "",
    accent_idx: int = 0,
    has_tpl: bool = False,
) -> None:
    """Render a KPI metric card with shadow, gradient, and prominent value.

    Matches the sample outputs: big number, short label, subtle description,
    shadow + gradient for visual depth.
    """
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)

    # Style: gradient + shadow + rounded corners
    if has_tpl:
        style_card(
            card,
            theme_color=_ACCENT_THEME[accent_idx % len(_ACCENT_THEME)],
            shadow_preset="card",
            corner_radius=8000,
        )
    else:
        c1 = _hex_to_rgb(_FALLBACK_HEX[accent_idx % len(_FALLBACK_HEX)])
        # Derive darker shade for gradient stop 2
        c2 = RGBColor(
            max(c1[0] - 30, 0), max(c1[1] - 30, 0), max(c1[2] - 30, 0)
        )
        style_card(
            card,
            gradient_stops=[(0.0, c1), (1.0, c2)],
            shadow_preset="card",
            corner_radius=8000,
        )

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(100000)
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    _set_autofit(tf)

    # Hero value
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Label
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

    # Description
    if description:
        p3 = tf.add_paragraph()
        p3.text = description[:80]
        p3.font.size = Pt(8)
        p3.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        p3.alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# Content Card  —  header strip + white body + items + shadow
# ---------------------------------------------------------------------------

def render_content_card(
    slide,
    x: int, y: int, w: int, h: int,
    title: str,
    items: list[str],
    accent_idx: int = 0,
    has_tpl: bool = False,
) -> None:
    """Card with coloured header strip, light body, shadow, and bullet items."""
    # Light background card with shadow
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    if has_tpl:
        card.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        card.fill.fore_color.brightness = 0.96
    else:
        card.fill.fore_color.rgb = _hex_to_rgb("F7F9FC")
    remove_outline(card)
    add_shadow(card, preset="subtle")
    set_corner_radius(card, 8000)

    # Header strip
    hdr_h = Emu(320000)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, hdr_h)
    hdr.fill.solid()
    if has_tpl:
        hdr.fill.fore_color.theme_color = _ACCENT_THEME[accent_idx % len(_ACCENT_THEME)]
    else:
        hdr.fill.fore_color.rgb = _hex_to_rgb(_FALLBACK_HEX[accent_idx % len(_FALLBACK_HEX)])
    remove_outline(hdr)

    # Header title
    htf = hdr.text_frame
    htf.word_wrap = True
    htf.margin_left = Emu(80000)
    hp = htf.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(11)
    hp.font.bold = True
    hp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    hp.alignment = PP_ALIGN.LEFT

    # Body items
    if items:
        body_top = y + hdr_h + Emu(60000)
        body_h = h - hdr_h - Emu(120000)
        tx = slide.shapes.add_textbox(
            x + Emu(100000), body_top,
            w - Emu(200000), body_h
        )
        tf = tx.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.LEFT
            if idx > 0:
                p.space_before = Pt(4)
        _set_autofit(tf)


# ---------------------------------------------------------------------------
# Stat Callout  —  accent circle + number + label
# ---------------------------------------------------------------------------

def render_stat_callout(
    slide,
    x: int, y: int, w: int, h: int,
    number: str,
    label: str,
    accent_idx: int = 0,
    has_tpl: bool = False,
) -> None:
    """Small stat element: coloured circle with number + label below."""
    circle_d = min(w, Emu(500000))
    cx = x + (w - circle_d) // 2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, y, circle_d, circle_d)
    if has_tpl:
        style_card(
            circle,
            theme_color=_ACCENT_THEME[accent_idx % len(_ACCENT_THEME)],
            shadow_preset="subtle",
            corner_radius=None,  # ovals don't have corners
        )
    else:
        c = _hex_to_rgb(_FALLBACK_HEX[accent_idx % len(_FALLBACK_HEX)])
        c2 = RGBColor(max(c[0] - 20, 0), max(c[1] - 20, 0), max(c[2] - 20, 0))
        style_card(circle, gradient_stops=[(0.0, c), (1.0, c2)],
                   shadow_preset="subtle", corner_radius=None)

    # Number inside circle
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Label below
    label_y = y + circle_d + Emu(30000)
    label_h = h - circle_d - Emu(30000)
    if label_h > 0:
        tx = slide.shapes.add_textbox(x, label_y, w, label_h)
        ltf = tx.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(9)
        lp.alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# Accent Divider Bar  —  horizontal gradient accent line
# ---------------------------------------------------------------------------

def render_accent_divider(
    slide,
    x: int, y: int, w: int,
    has_tpl: bool = False,
    accent_idx: int = 0,
    height: int = Emu(40000),
) -> None:
    """Thin horizontal accent bar with gradient."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, height)
    if has_tpl:
        style_accent_bar(bar, theme_color=_ACCENT_THEME[accent_idx % len(_ACCENT_THEME)])
    else:
        c1 = _hex_to_rgb(_FALLBACK_HEX[accent_idx % len(_FALLBACK_HEX)])
        c2_idx = (accent_idx + 1) % len(_FALLBACK_HEX)
        c2 = _hex_to_rgb(_FALLBACK_HEX[c2_idx])
        style_accent_bar(bar, gradient_stops=[(0.0, c1), (1.0, c2)])


# ---------------------------------------------------------------------------
# Numbered Step Circle  —  like the 01/02/03 in the sample UAE output
# ---------------------------------------------------------------------------

def render_numbered_circle(
    slide,
    x: int, y: int, size: int,
    number: int,
    accent_idx: int = 0,
    has_tpl: bool = False,
) -> None:
    """Small numbered circle with gradient and shadow (matches UAE sample)."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    if has_tpl:
        style_numbered_circle(circle, _ACCENT_THEME[accent_idx % len(_ACCENT_THEME)])
    else:
        c = _hex_to_rgb(_FALLBACK_HEX[accent_idx % len(_FALLBACK_HEX)])
        c2 = RGBColor(max(c[0] - 15, 0), max(c[1] - 15, 0), max(c[2] - 15, 0))
        add_gradient(circle, [(0.0, c), (1.0, c2)], angle=270.0)
        add_shadow(circle, preset="subtle")
        remove_outline(circle)

    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{number:02d}"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# Chart Container  —  background card with shadow behind a chart area
# ---------------------------------------------------------------------------

def render_chart_container(
    slide,
    x: int, y: int, w: int, h: int,
    has_tpl: bool = False,
) -> None:
    """Add a subtle shadow/surface background behind a chart or table area.

    Charts (GraphicsFrame) can't receive shadows directly, so we place
    a rounded rect behind them.
    """
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x - Emu(30000), y - Emu(30000),
        w + Emu(60000), h + Emu(60000),
    )
    bg.fill.solid()
    if has_tpl:
        bg.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        bg.fill.fore_color.brightness = 0.97
    else:
        bg.fill.fore_color.rgb = _hex_to_rgb("FBFCFE")
    remove_outline(bg)
    add_shadow(bg, preset="subtle")
    set_corner_radius(bg, 6000)


# ---------------------------------------------------------------------------
# Section Badge  —  accent pill with section number for divider slides
# ---------------------------------------------------------------------------

def render_section_badge(
    slide,
    x: int, y: int,
    section_text: str,
    has_tpl: bool = False,
    accent_idx: int = 0,
) -> None:
    """Small coloured pill/badge for section dividers."""
    badge_w = Emu(1200000)
    badge_h = Emu(350000)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, badge_w, badge_h)
    if has_tpl:
        style_card(
            pill,
            theme_color=_ACCENT_THEME[accent_idx % len(_ACCENT_THEME)],
            shadow_preset="subtle",
            corner_radius=16000,
        )
    else:
        c = _hex_to_rgb(_FALLBACK_HEX[accent_idx % len(_FALLBACK_HEX)])
        c2 = RGBColor(max(c[0] - 20, 0), max(c[1] - 20, 0), max(c[2] - 20, 0))
        style_card(pill, gradient_stops=[(0.0, c), (1.0, c2)],
                   shadow_preset="subtle", corner_radius=16000)

    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = section_text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
