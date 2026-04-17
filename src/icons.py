"""Programmatic MSO-shape icon library.

Each icon is composed of 2-6 native PowerPoint auto-shapes (rectangles,
ovals, triangles, connectors) arranged inside an invisible bounding box.
Renders natively in PowerPoint / LibreOffice / Google Slides — no external
image files, no copyrighted assets.

Public API:
- ``draw_icon(slide, name, left, top, size, master_info, has_tpl)`` — draws
  the named icon at the given position/size.
- ``icon_for_keyword(text)`` — returns the best-matching icon name given a
  short label (used by ``icon_row_list`` archetype).
- ``ICON_NAMES`` — tuple of every supported icon name.
"""
from __future__ import annotations

import logging
import re
from typing import Callable

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from .schemas import SlideMasterInfo
from .color_utils import FALLBACK_ACCENT_HEX

logger = logging.getLogger(__name__)


# ── Helpers ─────────────────────────────────────────────────────────

def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _accent_hex(master_info: SlideMasterInfo | None) -> str:
    """Return the first accent hex with ≥3:1 contrast against a white slide.

    Falls back to the first accent, or to a safe corporate blue if none are
    defined. This guarantees icons remain visible on templates whose
    ACCENT_1 is a very light tint (e.g. UAE Solar ``EFF3E5``).
    """
    if master_info and master_info.theme_colors.accents():
        accents = master_info.theme_colors.accents()
        # Use the color_utils contrast helper lazily (avoid import at module load).
        from .color_utils import contrast_ratio
        for hex_val in accents:
            if hex_val and contrast_ratio(hex_val, "FFFFFF") >= 3.0:
                return hex_val.lstrip("#")
        return accents[0]
    return FALLBACK_ACCENT_HEX[0]


def _apply_accent(shape, master_info: SlideMasterInfo | None, has_tpl: bool,
                  *, outline: bool = False) -> None:
    """Fill (or outline) a shape with a readable accent colour.

    Explicitly resolves the hex (rather than relying on ``theme_color``) so
    pastel ACCENT_1 values don't render the icon as near-invisible on white.
    """
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(_accent_hex(master_info))
    if outline:
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()


def _apply_white(shape) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.fill.background()


# ── Icon primitives ─────────────────────────────────────────────────
# Each function receives (slide, left, top, size, master_info, has_tpl).
# *size* is total icon bounding-box edge in EMU.

def _icon_globe(slide, left, top, size, master_info, has_tpl):
    # Circle + two meridians + equator
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    _apply_accent(outer, master_info, has_tpl)
    # Inner circle (thin ring effect)
    inset = size // 8
    inner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + inset, top + inset, size - 2 * inset, size - 2 * inset
    )
    _apply_white(inner)
    # Equator + meridian lines
    mid_y = top + size // 2
    eq = slide.shapes.add_connector(1, left + inset, mid_y, left + size - inset, mid_y)
    eq.line.color.rgb = _hex_to_rgb(_accent_hex(master_info))
    eq.line.width = Pt(1.25)
    mid_x = left + size // 2
    mer = slide.shapes.add_connector(1, mid_x, top + inset, mid_x, top + size - inset)
    mer.line.color.rgb = _hex_to_rgb(_accent_hex(master_info))
    mer.line.width = Pt(1.25)
    # Oblique meridian (ellipse)
    obl = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + size // 3, top + inset,
        size // 3, size - 2 * inset
    )
    obl.fill.background()
    obl.line.color.rgb = _hex_to_rgb(_accent_hex(master_info))
    obl.line.width = Pt(1.25)


def _icon_people(slide, left, top, size, master_info, has_tpl):
    # Head (small circle) + body (rounded rect — ROUND_2_SAME_RECTANGLE rounds
    # the top-two corners by default, mimicking a shoulder silhouette).
    head_d = size // 3
    head = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + (size - head_d) // 2, top + size // 12,
        head_d, head_d,
    )
    _apply_accent(head, master_info, has_tpl)
    body_w = size * 3 // 4
    body_h = size // 2
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
        left + (size - body_w) // 2, top + size // 2 - size // 16,
        body_w, body_h,
    )
    _apply_accent(body, master_info, has_tpl)


def _icon_people_group(slide, left, top, size, master_info, has_tpl):
    # Three overlapping figures: centre (big), left (small), right (small)
    _icon_people(slide, left + size // 5, top + size // 10, size * 3 // 4, master_info, has_tpl)
    # Small silhouette left
    small = size // 2
    _icon_people(slide, left, top + size // 3, small, master_info, has_tpl)
    _icon_people(slide, left + size - small, top + size // 3, small, master_info, has_tpl)


def _icon_chart_bar(slide, left, top, size, master_info, has_tpl):
    # Three bars of increasing height inside a square frame
    bar_gap = size // 16
    bar_w = (size - 4 * bar_gap) // 3
    base_y = top + size - size // 10
    heights = [size // 3, size * 2 // 3 - size // 20, size - size // 5]
    for i, h in enumerate(heights):
        x = left + bar_gap + i * (bar_w + bar_gap)
        y = base_y - h
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, h)
        _apply_accent(bar, master_info, has_tpl)
    # Axis
    axis = slide.shapes.add_connector(1, left, base_y, left + size, base_y)
    axis.line.color.rgb = _hex_to_rgb(_accent_hex(master_info))
    axis.line.width = Pt(1.5)


def _icon_chart_up(slide, left, top, size, master_info, has_tpl):
    # Trend arrow going up-right + axis (rotate RIGHT_ARROW by -45° since
    # UP_RIGHT_ARROW is not exposed by python-pptx).
    base_y = top + size - size // 10
    axis = slide.shapes.add_connector(1, left, base_y, left + size, base_y)
    axis.line.color.rgb = _hex_to_rgb(_accent_hex(master_info))
    axis.line.width = Pt(1.5)
    arrow_len = int(size * 0.72)
    arrow_h = size // 3
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left + size // 2 - arrow_len // 2,
        top + size // 2 - arrow_h // 2,
        arrow_len, arrow_h,
    )
    _apply_accent(arrow, master_info, has_tpl)
    arrow.rotation = -40.0


def _icon_chart_down(slide, left, top, size, master_info, has_tpl):
    base_y = top + size - size // 10
    axis = slide.shapes.add_connector(1, left, base_y, left + size, base_y)
    axis.line.color.rgb = _hex_to_rgb(_accent_hex(master_info))
    axis.line.width = Pt(1.5)
    # Down-right arrow: rotate RIGHT_ARROW by +40°
    arrow_len = int(size * 0.72)
    arrow_h = size // 3
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left + size // 2 - arrow_len // 2,
        top + size // 2 - arrow_h // 2,
        arrow_len, arrow_h,
    )
    _apply_accent(arrow, master_info, has_tpl)
    arrow.rotation = 40.0


def _icon_chart_pie(slide, left, top, size, master_info, has_tpl):
    # Outer ring
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    _apply_white(outer)
    outer.line.color.rgb = _hex_to_rgb(_accent_hex(master_info))
    outer.line.width = Pt(1.5)
    # Filled pie slice (~30% of the pie) — pie shape
    pie = slide.shapes.add_shape(MSO_SHAPE.PIE, left, top, size, size)
    _apply_accent(pie, master_info, has_tpl)
    pie.adjustments[0] = 90 * 60000  # start angle
    pie.adjustments[1] = 210 * 60000  # end angle


def _icon_warning(slide, left, top, size, master_info, has_tpl):
    # Triangle with exclamation mark (correct spelling: ISOSCELES)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, size, size)
    _apply_accent(tri, master_info, has_tpl)
    # Exclamation mark (stem + dot)
    stem_w = size // 10
    stem_h = size // 3
    stem_x = left + (size - stem_w) // 2
    stem_y = top + size // 3
    stem = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, stem_x, stem_y, stem_w, stem_h)
    _apply_white(stem)
    dot_d = size // 8
    dot_x = left + (size - dot_d) // 2
    dot_y = top + size * 3 // 4
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, dot_d, dot_d)
    _apply_white(dot)


def _icon_lightbulb(slide, left, top, size, master_info, has_tpl):
    # Bulb (oval) + base (small rect)
    bulb_h = size * 2 // 3
    bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + size // 8, top, size * 3 // 4, bulb_h)
    _apply_accent(bulb, master_info, has_tpl)
    # Base
    base_w = size // 2
    base_h = size // 6
    base = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + (size - base_w) // 2, top + bulb_h - size // 20,
        base_w, base_h,
    )
    _apply_accent(base, master_info, has_tpl)
    # Filament
    fil_w = size // 6
    fil_h = size // 6
    fil = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + (size - fil_w) // 2, top + bulb_h - size // 20 + base_h,
        fil_w, fil_h,
    )
    _apply_accent(fil, master_info, has_tpl)


def _icon_shield(slide, left, top, size, master_info, has_tpl):
    # Shield: pentagon-ish shape approximated with rounded rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, size, size)
    _apply_accent(shape, master_info, has_tpl)
    shape.rotation = 180.0  # point down
    # Check mark inside
    cm_w = size // 2
    cm_h = size // 4
    cm = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left + (size - cm_w) // 2, top + (size - cm_h) // 2,
        cm_w, cm_h,
    )
    _apply_white(cm)


def _icon_gear(slide, left, top, size, master_info, has_tpl):
    # Cross-shaped gear approximation: outer square rotated + inner circle
    outer = slide.shapes.add_shape(MSO_SHAPE.SUN, left, top, size, size)
    _apply_accent(outer, master_info, has_tpl)
    # Central white hub
    hub_d = size // 3
    hub = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + (size - hub_d) // 2, top + (size - hub_d) // 2,
        hub_d, hub_d,
    )
    _apply_white(hub)


def _icon_clock(slide, left, top, size, master_info, has_tpl):
    face = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    _apply_accent(face, master_info, has_tpl)
    inner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + size // 12, top + size // 12,
        size * 5 // 6, size * 5 // 6,
    )
    _apply_white(inner)
    # Hour hand (vertical) + minute hand (horizontal)
    cx = left + size // 2
    cy = top + size // 2
    accent = _hex_to_rgb(_accent_hex(master_info))
    hour = slide.shapes.add_connector(1, cx, cy, cx, cy - size // 3)
    hour.line.color.rgb = accent
    hour.line.width = Pt(2.0)
    minute = slide.shapes.add_connector(1, cx, cy, cx + size // 3, cy)
    minute.line.color.rgb = accent
    minute.line.width = Pt(1.5)


def _icon_money(slide, left, top, size, master_info, has_tpl):
    # Accent-filled coin overlaid with a bold white "$" glyph. Using a text
    # box gives a crisp, unambiguous dollar sign — much clearer than trying
    # to build the S-curve out of primitive shapes.
    coin = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    _apply_accent(coin, master_info, has_tpl)

    # Dollar-sign textbox, sized to roughly the coin's inner area.
    glyph = slide.shapes.add_textbox(left, top, size, size)
    tf = glyph.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    try:
        from pptx.enum.text import MSO_ANCHOR as _A
        tf.vertical_anchor = _A.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "$"
    # Size the glyph so it nearly fills the coin (Emu → pt conversion: 12700 EMU/pt)
    font_pt = max(int(size / 15000), 12)
    p.font.size = Pt(font_pt)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _icon_location(slide, left, top, size, master_info, has_tpl):
    # Teardrop = circle on top + inverted isosceles triangle below
    drop = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + size // 8, top, size * 3 // 4, size * 3 // 4)
    _apply_accent(drop, master_info, has_tpl)
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        left + size // 4, top + size // 2,
        size // 2, size // 2,
    )
    _apply_accent(tri, master_info, has_tpl)
    tri.rotation = 180.0
    # Inner white dot
    dot_d = size // 4
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + (size - dot_d) // 2, top + size // 4,
        dot_d, dot_d,
    )
    _apply_white(dot)


def _icon_document(slide, left, top, size, master_info, has_tpl):
    # Folded-corner rectangle with three text lines
    page_w = size * 3 // 4
    page_h = size * 9 // 10
    page = slide.shapes.add_shape(
        MSO_SHAPE.FOLDED_CORNER,
        left + size // 8, top + size // 20,
        page_w, page_h,
    )
    _apply_accent(page, master_info, has_tpl)
    # Text lines (white horizontal bars)
    line_w = page_w * 2 // 3
    line_h = page_h // 14
    for i in range(3):
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + size // 8 + (page_w - line_w) // 2,
            top + size // 20 + page_h // 3 + i * (line_h * 2 + line_h // 2),
            line_w, line_h,
        )
        _apply_white(ln)


def _icon_check(slide, left, top, size, master_info, has_tpl):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    _apply_accent(circle, master_info, has_tpl)
    # Check mark (right-arrow rotated)
    cm = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left + size // 5, top + size // 3,
        size * 3 // 5, size // 4,
    )
    _apply_white(cm)


def _icon_x(slide, left, top, size, master_info, has_tpl):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    _apply_accent(circle, master_info, has_tpl)
    # Two crossing rectangles forming an X
    bar_w = size // 10
    bar_len = size * 3 // 5
    center_x = left + size // 2
    center_y = top + size // 2
    bar1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        center_x - bar_len // 2, center_y - bar_w // 2,
        bar_len, bar_w,
    )
    _apply_white(bar1)
    bar1.rotation = 45.0
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        center_x - bar_len // 2, center_y - bar_w // 2,
        bar_len, bar_w,
    )
    _apply_white(bar2)
    bar2.rotation = -45.0


def _icon_flag(slide, left, top, size, master_info, has_tpl):
    # Pole + flag
    pole_w = size // 16
    pole = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + size // 8, top, pole_w, size,
    )
    _apply_accent(pole, master_info, has_tpl)
    # Flag (triangle pennant)
    flag = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        left + size // 8 + pole_w, top + size // 10,
        size * 2 // 3, size // 3,
    )
    _apply_accent(flag, master_info, has_tpl)


def _icon_star(slide, left, top, size, master_info, has_tpl):
    star = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, left, top, size, size)
    _apply_accent(star, master_info, has_tpl)


# ── Registry + keyword mapping ──────────────────────────────────────

_ICON_REGISTRY: dict[str, Callable] = {
    "globe": _icon_globe,
    "people": _icon_people,
    "people_group": _icon_people_group,
    "chart_bar": _icon_chart_bar,
    "chart_up": _icon_chart_up,
    "chart_down": _icon_chart_down,
    "chart_pie": _icon_chart_pie,
    "warning": _icon_warning,
    "lightbulb": _icon_lightbulb,
    "shield": _icon_shield,
    "gear": _icon_gear,
    "clock": _icon_clock,
    "money": _icon_money,
    "location": _icon_location,
    "document": _icon_document,
    "check": _icon_check,
    "x_mark": _icon_x,
    "flag": _icon_flag,
    "star": _icon_star,
}

ICON_NAMES: tuple[str, ...] = tuple(_ICON_REGISTRY.keys())


# Keyword → icon mappings. We use *word-start* anchors (``\b`` on the left only)
# so that "challenges" still matches "challenge", "growing" still matches "grow", etc.
_KEYWORD_MAP: list[tuple[str, str]] = [
    # Geography / markets
    (r"\b(global|world|international|geograph|countr|market|export|trade)", "globe"),
    (r"\b(location|region|city|site|plant|hq|office|headquart)", "location"),
    (r"\b(flag|national|sovereign)", "flag"),
    # People / talent / workforce
    (r"\b(team|talent|workforce|staff|employee|customer|investor|user|people)", "people_group"),
    (r"\b(presenter|speaker|executive|leader|founder|ceo|cfo)", "people"),
    # Charts / data
    (r"\b(growth|grow|increas|rise|gain|expan|scale|surge|accelerat|upsid|boost)", "chart_up"),
    (r"\b(declin|decreas|drop|fall|downsid|contract|shrink|slowdown|reduc)", "chart_down"),
    (r"\b(share|split|portion|proportion|allocation|breakdown|distribution)", "chart_pie"),
    (r"\b(benchmark|comparison|metric|kpi|statistic|data|indicator|figure|chart)", "chart_bar"),
    # Risk / warning / challenges
    (r"\b(risk|threat|warning|alert|danger|concern|issue|problem|challenge|headwind|vulnerab)", "warning"),
    # Ideas / innovation / strategy / recommendations
    (r"\b(idea|insight|innovat|strategy|strategi|vision|concept|approach|method|recommend|propos)", "lightbulb"),
    # Security / governance / compliance
    (r"\b(security|compliance|regulat|governance|safeguard|protect|shield|defense|defence|privacy|policy)", "shield"),
    # Process / operations / tech
    (r"\b(process|operation|workflow|pipeline|system|mechanism|engine|technolog|automat|manufactur)", "gear"),
    # Time / timeline / milestones
    (r"\b(timeline|duration|schedule|milestone|roadmap|year|quarter|period|history|chronolog|phase)", "clock"),
    # Money / finance / valuation
    (r"\b(revenue|cost|profit|invest|capital|money|price|value|valuation|financ|economic|spend|budget|roi)", "money"),
    # Reports / documentation
    (r"\b(report|document|framework|disclosure|publication|research|analys|stud|review)", "document"),
    # Outcomes — positive
    (r"\b(success|achievement|complete|approved|win|positive|improve|deliver|accomplish)", "check"),
    # Outcomes — negative
    (r"\b(failure|reject|miss|deni|negative|blocker)", "x_mark"),
    # Highlights
    (r"\b(highlight|featured|key|primary|hero|important|priority|critical|signature)", "star"),
]


def icon_for_keyword(text: str, fallback: str = "chart_bar") -> str:
    """Return the best-matching icon name for a short label.

    Matches are case-insensitive and performed against whole-word regex
    patterns. Falls back to ``fallback`` if no keyword matches.
    """
    if not text:
        return fallback
    lower = text.lower()
    for pattern, name in _KEYWORD_MAP:
        if re.search(pattern, lower):
            return name
    return fallback


def draw_icon(slide, name: str, left, top, size: int,
              master_info: SlideMasterInfo | None = None,
              has_tpl: bool = False) -> None:
    """Draw the named icon onto *slide* at the given position.

    *size* is the bounding-box edge length in EMU (square icons).
    Unknown names silently fall back to ``chart_bar``.
    """
    fn = _ICON_REGISTRY.get(name) or _ICON_REGISTRY["chart_bar"]
    try:
        fn(slide, int(left), int(top), int(size), master_info, has_tpl)
    except Exception as exc:
        # Icon rendering errors should never break the slide, but they MUST be
        # logged — silently dropping icons is exactly what caused the empty-row
        # defect in earlier builds. Fall back to a simple accent circle so the
        # row still has a visual anchor.
        logger.warning("Icon %r failed to render (%s); using fallback circle", name, exc)
        try:
            fallback = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, int(left), int(top), int(size), int(size)
            )
            _apply_accent(fallback, master_info, has_tpl)
        except Exception:
            pass
