"""Slide-master introspection and layout selection.

Critical design decisions (from research §6.2):
- Layout selection is **index-based**, never name-based, because UAE Solar
  has 3 layouts with the identical name "1_E_Title, Subtitle and Body".
- UAE Solar has **no Blank layout**; the emptiest layout (index 4) has a
  full-width image covering 61 % of the slide.  The grid system must
  account for this by selecting the layout with the fewest content
  placeholders as the canvas layout.
- Template identification uses fuzzy word-overlap scoring against the
  markdown filename to auto-detect which of the 3 templates to use.
"""
from __future__ import annotations
import zipfile
import logging
from pathlib import Path
from lxml import etree
from pptx import Presentation
from .schemas import SlideMasterInfo, LayoutInfo, PlaceholderInfo, ThemeColors
from . import config

logger = logging.getLogger(__name__)

_DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS = {"a": _DRAWINGML_NS}

_COLOR_SLOTS = (
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
)

# ── Per-template index-based layout map (research §6.2) ─────────────
# Keys: purpose strings.  Values: layout indices.
# This avoids name-based lookup entirely — UAE Solar has 3 layouts with
# the same name, so name-matching is unreliable.
LAYOUT_MAP: dict[str, dict[str, int]] = {
    "AI_Bubble": {"cover": 0, "divider": 1, "blank": 2, "title_only": 3, "end": 4},
    "UAE_Solar": {"cover": 0, "divider": 1, "content": 2, "end": -1},
    "Accenture": {"cover": 0, "divider": 2, "blank": 3, "title_only": 4, "end": 5},
}


def _extract_theme_colors(template_path: Path) -> ThemeColors:
    """Parse the PPTX theme XML and return all 12 color-scheme slots."""
    colors: dict[str, str] = {}
    try:
        with zipfile.ZipFile(str(template_path), "r") as z:
            theme_files = sorted(
                f for f in z.namelist()
                if f.startswith("ppt/theme/theme") and f.endswith(".xml")
            )
            if not theme_files:
                return ThemeColors()
            theme_xml = z.read(theme_files[0])

        root = etree.fromstring(theme_xml)
        clr_scheme = root.find(".//a:clrScheme", _NS)
        if clr_scheme is None:
            return ThemeColors()

        for child in clr_scheme:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag not in _COLOR_SLOTS:
                continue
            srgb = child.find("a:srgbClr", _NS)
            sys_clr = child.find("a:sysClr", _NS)
            if srgb is not None:
                colors[tag] = srgb.get("val", "000000")
            elif sys_clr is not None:
                colors[tag] = sys_clr.get("lastClr", "000000")

    except Exception as e:
        logger.warning(f"Failed to extract theme colors: {e}")

    return ThemeColors(**colors) if colors else ThemeColors()


def _categorize_layout(name: str, placeholders: list[PlaceholderInfo]) -> str:
    """Categorize a layout based on its name and placeholder structure."""
    name_lower = name.lower().strip()

    # ── Name-based heuristics (order matters) ──
    if any(kw in name_lower for kw in ("cover", "title company", "title slide")):
        return "cover"
    if any(kw in name_lower for kw in ("divider", "section", "header")):
        return "divider"
    if "thank" in name_lower or "end" in name_lower:
        return "thank_you"
    if name_lower == "blank" or (not placeholders and "blank" in name_lower):
        return "blank"
    if "title only" in name_lower:
        return "title_only"
    if "two content" in name_lower or "comparison" in name_lower:
        return "two_content"
    if "content" in name_lower or "title, content" in name_lower:
        return "title_content"

    # ── Placeholder-structure heuristics ──
    # Filter out footer / slide-number / date placeholders
    content_phs = [p for p in placeholders
                   if (p.ph_type or "").upper() not in ("SLIDE_NUMBER", "FOOTER", "DATE_TIME")]

    has_title = any((p.ph_type or "").upper() in ("TITLE", "CENTER_TITLE") for p in content_phs)
    has_center_title = any((p.ph_type or "").upper() == "CENTER_TITLE" for p in content_phs)
    body_count = sum(1 for p in content_phs if (p.ph_type or "").upper() in ("BODY", "OBJECT", "TABLE", "CHART"))

    if not content_phs:
        return "blank"
    if has_center_title and body_count == 0:
        return "cover"  # center title with no body → likely cover or divider
    if has_title and body_count >= 2:
        return "two_content"
    if has_title and body_count == 1:
        return "title_content"
    if has_title and body_count == 0:
        return "title_only"

    return "other"


def read_slide_master(template_path: str | Path) -> SlideMasterInfo:
    """Read a .pptx template and extract layouts, placeholders, and theme colors."""
    template_path = Path(template_path)
    prs = Presentation(str(template_path))

    theme_colors = _extract_theme_colors(template_path)
    logger.info(f"Theme accents: {theme_colors.accents()}")

    layouts: list[LayoutInfo] = []
    master = prs.slide_masters[0]

    for idx, layout in enumerate(master.slide_layouts):
        placeholders: list[PlaceholderInfo] = []
        for ph in layout.placeholders:
            fmt = ph.placeholder_format
            placeholders.append(PlaceholderInfo(
                idx=fmt.idx,
                name=ph.name,
                ph_type=str(fmt.type).split(".")[-1].split("(")[0].strip() if fmt.type else None,
                left=ph.left or 0,
                top=ph.top or 0,
                width=ph.width or 0,
                height=ph.height or 0,
            ))

        category = _categorize_layout(layout.name, placeholders)
        layouts.append(LayoutInfo(
            index=idx,
            name=layout.name,
            category=category,
            placeholders=placeholders,
        ))
        logger.debug(f"Layout [{idx}] '{layout.name}' → {category} ({len(placeholders)} phs)")

    return SlideMasterInfo(
        template_path=str(template_path),
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        layouts=layouts,
        theme_colors=theme_colors,
    )


def identify_template(template_path: str | Path) -> str | None:
    """Identify which known template family a PPTX belongs to.

    Returns one of the LAYOUT_MAP keys ("AI_Bubble", "UAE_Solar",
    "Accenture") or *None* if unrecognised.
    """
    stem = Path(template_path).stem.lower()
    if "ai_bubble" in stem or "ai bubble" in stem or "detection" in stem and "investment" in stem:
        return "AI_Bubble"
    if "uae" in stem or "solar" in stem and "2050" in stem:
        return "UAE_Solar"
    if "accenture" in stem or "acquisition" in stem:
        return "Accenture"
    return None


def find_layout_by_category(master_info: SlideMasterInfo, category: str,
                            excluded_idx: int | None = None) -> LayoutInfo | None:
    """Find the first layout matching a given category, optionally skipping an index."""
    for layout in master_info.layouts:
        if layout.category == category and layout.index != excluded_idx:
            return layout
    return None


def _find_emptiest_layout(master_info: SlideMasterInfo,
                          excluded_idx: int | None = None) -> LayoutInfo:
    """Find the layout with the fewest content placeholders (canvas layout).

    This handles templates like UAE Solar that have no Blank layout.
    The "emptiest" layout is the best canvas for programmatic content.
    """
    _FURNITURE_TYPES = frozenset({"SLIDE_NUMBER", "FOOTER", "DATE_TIME"})
    best: LayoutInfo | None = None
    best_count = 999
    for layout in master_info.layouts:
        if layout.index == excluded_idx:
            continue
        if layout.category in ("cover", "divider", "thank_you"):
            continue
        content_phs = [
            p for p in layout.placeholders
            if (p.ph_type or "").upper() not in _FURNITURE_TYPES
        ]
        if len(content_phs) < best_count:
            best_count = len(content_phs)
            best = layout
    return best or master_info.layouts[-1]


def get_layout_for_slide_type(master_info: SlideMasterInfo, slide_type: str,
                              excluded_idx: int | None = None) -> LayoutInfo:
    """Map a slide_type to the best available layout.

    Strategy (research §6.2–6.3):
    1. If the template is known (in LAYOUT_MAP), use the explicit index.
    2. Otherwise fall back to category-based heuristic matching.
    3. If no category matches (e.g. UAE Solar has no blank), select the
       layout with the fewest content placeholders as the canvas.

    *excluded_idx* prevents the bookend closing layout from being picked
    for content slides.
    """
    # ── Step 1: Try known template index map ──
    tpl_name = identify_template(master_info.template_path)
    if tpl_name and tpl_name in LAYOUT_MAP:
        tpl_map = LAYOUT_MAP[tpl_name]
        # Map slide_type → purpose key in the template map
        purpose_map: dict[str, list[str]] = {
            "cover": ["cover"],
            "section_divider": ["divider"],
            "thank_you": ["end", "cover"],
            "agenda": ["title_only", "blank", "content"],
            "executive_summary": ["title_only", "blank", "content"],
            "content": ["title_only", "blank", "content"],
            "chart": ["title_only", "blank", "content"],
            "table": ["title_only", "blank", "content"],
            "infographic": ["title_only", "blank", "content"],
            "mixed": ["title_only", "blank", "content"],
            "conclusion": ["title_only", "blank", "content"],
        }
        for purpose in purpose_map.get(slide_type, ["blank", "content"]):
            idx = tpl_map.get(purpose, -1)
            if idx >= 0 and idx != excluded_idx and idx < len(master_info.layouts):
                return master_info.layouts[idx]

    # ── Step 2: Category-based fallback ──
    category_map: dict[str, list[str]] = {
        "cover": ["cover"],
        "section_divider": ["divider", "cover"],
        "thank_you": ["thank_you", "cover", "title_content"],
        "agenda": ["title_only", "title_content", "blank"],
        "executive_summary": ["title_only", "title_content", "blank"],
        "content": ["title_only", "title_content", "blank"],
        "chart": ["title_only", "blank"],
        "table": ["title_only", "blank"],
        "infographic": ["title_only", "blank"],
        "mixed": ["title_only", "two_content", "blank"],
        "conclusion": ["title_only", "title_content", "blank"],
    }

    candidates = category_map.get(slide_type, ["blank"])
    for cat in candidates:
        layout = find_layout_by_category(master_info, cat, excluded_idx=excluded_idx)
        if layout is not None:
            return layout

    # ── Step 3: Emptiest-layout fallback (handles UAE Solar no-blank) ──
    return _find_emptiest_layout(master_info, excluded_idx=excluded_idx)


def auto_detect_template(md_filename: str) -> Path | None:
    """Find the best matching template using fuzzy word-overlap scoring."""
    import re
    templates_dir = config.TEMPLATES_DIR
    if templates_dir is None or not templates_dir.exists():
        return None

    # Normalize filename into keyword set
    md_stem = Path(md_filename).stem.lower()
    md_words = set(re.findall(r'[a-z]{3,}', md_stem))

    best_match = None
    best_score = 0

    for tpl in templates_dir.glob("*.pptx"):
        tpl_stem = tpl.stem.lower().replace("template_", "")
        tpl_words = set(re.findall(r'[a-z]{3,}', tpl_stem))

        # Word overlap score
        overlap = len(md_words & tpl_words)
        # Substring bonus
        if tpl_stem in md_stem or md_stem in tpl_stem:
            overlap += 5

        if overlap > best_score:
            best_score = overlap
            best_match = tpl

    if best_match:
        logger.info(f"Template match: {best_match.name} (score={best_score})")
    else:
        templates = list(templates_dir.glob("*.pptx"))
        if templates:
            best_match = templates[0]
            logger.info(f"No keyword match, using first template: {best_match.name}")

    return best_match
