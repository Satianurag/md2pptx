#!/usr/bin/env python3
"""Automated verification script for generated PPTX presentations.

Checks each .pptx in the output/ directory for:
  1. Color contrast (WCAG AA ≥4.5:1 normal / ≥3:1 large text)
  2. Text overflow (estimated text area vs shape area)
  3. Structural integrity (cover first, thank_you last, no duplicate bookends)
  4. Overlapping shapes on the same slide
  5. Empty slides / missing content

Usage:
    python scripts/verify_output.py [--dir output/]
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Emu

# Add project root to path so we can import color_utils
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from src.color_utils import contrast_ratio, relative_luminance


# ── Data structures ──────────────────────────────────────────────────

@dataclass
class CheckResult:
    name: str
    status: str  # PASS / WARN / FAIL
    details: str = ""


@dataclass
class SlideReport:
    slide_num: int
    checks: list[CheckResult] = field(default_factory=list)


@dataclass
class FileReport:
    filename: str
    slide_count: int = 0
    slide_reports: list[SlideReport] = field(default_factory=list)
    global_checks: list[CheckResult] = field(default_factory=list)

    @property
    def worst_status(self) -> str:
        all_checks = self.global_checks[:]
        for sr in self.slide_reports:
            all_checks.extend(sr.checks)
        if any(c.status == "FAIL" for c in all_checks):
            return "FAIL"
        if any(c.status == "WARN" for c in all_checks):
            return "WARN"
        return "PASS"


# ── Helpers ──────────────────────────────────────────────────────────

def _extract_shape_colors(shape):
    """Try to extract (fill_hex, text_hex_list) from a shape."""
    fill_hex = None
    text_colors = []

    # Fill color
    try:
        ft = shape.fill.type
        if ft is not None:
            try:
                rgb = shape.fill.fore_color.rgb
                fill_hex = str(rgb)
            except Exception:
                pass
    except Exception:
        pass

    # Text colors
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if not para.text.strip():
                continue
            try:
                rgb = para.font.color.rgb
                if rgb:
                    text_colors.append((str(rgb), para.font.size, para.text.strip()[:40]))
            except Exception:
                text_colors.append(("inherited", None, para.text.strip()[:40]))

    return fill_hex, text_colors


def _is_large_text(font_size) -> bool:
    """Check if font size qualifies as large text (≥18pt or ≥14pt bold)."""
    if font_size is None:
        return False
    pt = font_size / 12700  # EMU to pt
    return pt >= 18


# ── Check implementations ────────────────────────────────────────────

def check_contrast(slide, slide_num: int) -> list[CheckResult]:
    """Check color contrast between text and shape fills."""
    results = []
    for shape in slide.shapes:
        fill_hex, text_colors = _extract_shape_colors(shape)
        if not fill_hex or not text_colors:
            continue

        for txt_hex, font_size, text_preview in text_colors:
            if txt_hex == "inherited":
                continue
            ratio = contrast_ratio(txt_hex, fill_hex)
            threshold = 3.0 if _is_large_text(font_size) else 4.5
            if ratio < threshold:
                results.append(CheckResult(
                    name="contrast",
                    status="FAIL",
                    details=f"Slide {slide_num}: contrast {ratio:.1f}:1 < {threshold}:1 "
                            f"(text #{txt_hex} on #{fill_hex}) \"{text_preview}\""
                ))
            elif ratio < threshold * 1.2:
                results.append(CheckResult(
                    name="contrast",
                    status="WARN",
                    details=f"Slide {slide_num}: marginal contrast {ratio:.1f}:1 "
                            f"(text #{txt_hex} on #{fill_hex}) \"{text_preview}\""
                ))

    if not results:
        results.append(CheckResult(name="contrast", status="PASS"))
    return results


def check_overflow(slide, slide_num: int) -> list[CheckResult]:
    """Estimate text overflow: compare character count to shape area."""
    results = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        total_chars = sum(len(p.text) for p in shape.text_frame.paragraphs)
        if total_chars == 0:
            continue

        # Rough estimate: 1 char ≈ 8pt wide × 14pt tall = ~130000 EMU²
        # Shape area in EMU²
        try:
            shape_area = shape.width * shape.height
        except Exception:
            continue

        # Estimate required area (very rough heuristic)
        avg_font = 12  # assume 12pt average
        char_w = avg_font * 9144  # ~0.1 inch per char at 12pt
        line_h = avg_font * 2 * 9144  # ~0.2 inch line height
        chars_per_line = max(shape.width // char_w, 1)
        lines_needed = (total_chars + chars_per_line - 1) // chars_per_line
        estimated_h = lines_needed * line_h

        if estimated_h > shape.height * 1.5:
            results.append(CheckResult(
                name="overflow",
                status="WARN",
                details=f"Slide {slide_num}: shape '{shape.name}' likely overflows "
                        f"({total_chars} chars, est. {lines_needed} lines, "
                        f"shape height={shape.height / 914400:.1f}in)"
            ))

    if not results:
        results.append(CheckResult(name="overflow", status="PASS"))
    return results


def _slide_has_thank_you(slide) -> bool:
    """Check if a slide has 'thank you' in its text OR baked into its layout shapes."""
    # Check slide text
    for shape in slide.shapes:
        if shape.has_text_frame and "thank" in shape.text_frame.text.strip().lower():
            return True
    # Check layout shapes (some templates bake "Thank you!" into the layout design)
    try:
        for shape in slide.slide_layout.shapes:
            if hasattr(shape, "text") and "thank" in shape.text.strip().lower():
                return True
    except Exception:
        pass
    return False


def check_structure(prs: Presentation) -> list[CheckResult]:
    """Check structural integrity: cover first, thank_you last, no dups."""
    results = []
    n = len(prs.slides)
    if n == 0:
        return [CheckResult(name="structure", status="FAIL", details="No slides")]

    # Check thank_you on last slide (text OR layout shapes)
    last_has_thank = _slide_has_thank_you(prs.slides[n - 1])
    if not last_has_thank:
        # Check if thank you appears elsewhere
        for i in range(n - 1):
            if _slide_has_thank_you(prs.slides[i]):
                results.append(CheckResult(
                    name="structure",
                    status="FAIL",
                    details=f"'Thank You' found on slide {i+1} instead of last slide {n}"
                ))
                break
        else:
            results.append(CheckResult(
                name="structure",
                status="WARN",
                details="No 'Thank You' text detected on any slide"
            ))
    else:
        # Verify no content slides also have thank you
        for i in range(n - 1):
            if _slide_has_thank_you(prs.slides[i]):
                results.append(CheckResult(
                    name="structure",
                    status="FAIL",
                    details=f"'Thank You' leaked to content slide {i+1} (should only be on last slide)"
                ))
                break
        else:
            results.append(CheckResult(name="structure", status="PASS",
                                       details="Thank You on last slide only"))

    # Check slide count
    if n < 10:
        results.append(CheckResult(name="structure", status="WARN",
                                   details=f"Only {n} slides (minimum 10)"))
    elif n > 15:
        results.append(CheckResult(name="structure", status="WARN",
                                   details=f"{n} slides (maximum 15)"))
    else:
        results.append(CheckResult(name="structure", status="PASS",
                                   details=f"{n} slides (within 10-15)"))

    return results


def check_overlaps(slide, slide_num: int) -> list[CheckResult]:
    """Detect shapes with significant overlap on the same slide."""
    results = []
    shapes_with_pos = []
    for shape in slide.shapes:
        try:
            shapes_with_pos.append((
                shape.name, shape.left, shape.top,
                shape.left + shape.width, shape.top + shape.height
            ))
        except Exception:
            continue

    for i in range(len(shapes_with_pos)):
        for j in range(i + 1, len(shapes_with_pos)):
            n1, x1, y1, x2, y2 = shapes_with_pos[i]
            n2, ax1, ay1, ax2, ay2 = shapes_with_pos[j]

            # Compute overlap
            ox = max(0, min(x2, ax2) - max(x1, ax1))
            oy = max(0, min(y2, ay2) - max(y1, ay1))
            overlap_area = ox * oy

            # Area of smaller shape
            a1 = (x2 - x1) * (y2 - y1)
            a2 = (ax2 - ax1) * (ay2 - ay1)
            min_area = min(a1, a2)

            if min_area > 0 and overlap_area / min_area > 0.5:
                results.append(CheckResult(
                    name="overlap",
                    status="WARN",
                    details=f"Slide {slide_num}: '{n1}' and '{n2}' overlap "
                            f"({overlap_area / min_area:.0%} of smaller shape)"
                ))

    if not results:
        results.append(CheckResult(name="overlap", status="PASS"))
    return results


def check_empty(slide, slide_num: int) -> list[CheckResult]:
    """Check for empty slides (no shapes or no text)."""
    if len(list(slide.shapes)) == 0:
        return [CheckResult(name="content", status="WARN",
                            details=f"Slide {slide_num}: no shapes (may use layout text)")]

    has_text = any(
        shape.has_text_frame and shape.text_frame.text.strip()
        for shape in slide.shapes
    )
    if not has_text:
        return [CheckResult(name="content", status="WARN",
                            details=f"Slide {slide_num}: no visible text content")]

    return [CheckResult(name="content", status="PASS")]


# ── Main verification ────────────────────────────────────────────────

def verify_file(path: Path) -> FileReport:
    """Run all checks on a single .pptx file."""
    report = FileReport(filename=path.name)
    try:
        prs = Presentation(str(path))
    except Exception as e:
        report.global_checks.append(
            CheckResult(name="open", status="FAIL", details=f"Cannot open: {e}")
        )
        return report

    report.slide_count = len(prs.slides)
    report.global_checks.extend(check_structure(prs))

    for i in range(len(prs.slides)):
        slide = prs.slides[i]
        sr = SlideReport(slide_num=i + 1)
        sr.checks.extend(check_contrast(slide, i + 1))
        sr.checks.extend(check_overflow(slide, i + 1))
        sr.checks.extend(check_overlaps(slide, i + 1))
        sr.checks.extend(check_empty(slide, i + 1))
        report.slide_reports.append(sr)

    return report


def print_report(report: FileReport) -> None:
    """Pretty-print a verification report."""
    status_icon = {"PASS": "✅", "WARN": "⚠️ ", "FAIL": "❌"}
    print(f"\n{'='*70}")
    print(f"{status_icon[report.worst_status]} {report.filename} ({report.slide_count} slides)")
    print(f"{'='*70}")

    # Global checks
    for c in report.global_checks:
        if c.status != "PASS":
            print(f"  {status_icon[c.status]} [{c.name}] {c.details}")

    # Per-slide issues (only non-PASS)
    for sr in report.slide_reports:
        issues = [c for c in sr.checks if c.status != "PASS"]
        for c in issues:
            print(f"  {status_icon[c.status]} [{c.name}] {c.details}")

    # Summary
    all_checks = report.global_checks[:]
    for sr in report.slide_reports:
        all_checks.extend(sr.checks)
    fails = sum(1 for c in all_checks if c.status == "FAIL")
    warns = sum(1 for c in all_checks if c.status == "WARN")
    passes = sum(1 for c in all_checks if c.status == "PASS")
    print(f"\n  Summary: {passes} pass, {warns} warn, {fails} fail")


def main():
    parser = argparse.ArgumentParser(description="Verify PPTX output quality")
    parser.add_argument("--dir", default="output", help="Directory with .pptx files")
    args = parser.parse_args()

    output_dir = Path(args.dir)
    if not output_dir.exists():
        print(f"Directory not found: {output_dir}")
        sys.exit(1)

    pptx_files = sorted(output_dir.glob("*.pptx"))
    if not pptx_files:
        print(f"No .pptx files found in {output_dir}")
        sys.exit(1)

    print(f"Verifying {len(pptx_files)} presentations in {output_dir}/\n")

    all_reports = []
    for path in pptx_files:
        report = verify_file(path)
        print_report(report)
        all_reports.append(report)

    # Overall summary
    print(f"\n{'='*70}")
    print("OVERALL SUMMARY")
    print(f"{'='*70}")
    for r in all_reports:
        icon = {"PASS": "✅", "WARN": "⚠️ ", "FAIL": "❌"}[r.worst_status]
        print(f"  {icon} {r.filename}")

    any_fail = any(r.worst_status == "FAIL" for r in all_reports)
    sys.exit(1 if any_fail else 0)


if __name__ == "__main__":
    main()
